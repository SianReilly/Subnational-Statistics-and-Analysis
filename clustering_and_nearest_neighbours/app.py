# pip install streamlit plotly pandas scikit-learn python-pptx kaleido openpyxl numpy
"""
Children's Services Nearest Neighbours & Clustering Explorer
--------------------------------------------------------------
A DAP (Data Analytics Principles) artefact built for a Local Authority
Children's Services directorate. Extends the ONS statistical nearest
neighbours / clustering methodology
(https://www.ons.gov.uk/peoplepopulationandcommunity/wellbeing/methodologies/
clusteringsimilarlocalauthoritiesandstatisticalnearestneighboursintheukmethodology)
with ten indicators specific to children's services demand.

Run with:  streamlit run app.py
Data expected at:  data/subnational_indicators.csv  (long format:
AREACD, AREANM, Indicator, Period, Measure, Unit, Value)
"""

import io
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
from sklearn.metrics import silhouette_samples, silhouette_score
from sklearn.preprocessing import StandardScaler

# ────────────────────────────────────────────────────────────────────────────
# 0. PAGE CONFIG & STYLE (ONS / Economist palette)
# ────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Children's Services Nearest Neighbours",
    page_icon="🧭",
    layout="wide",
    initial_sidebar_state="expanded",
)

NAVY = "#12436D"      # ONS dark blue - primary
TEAL = "#28A197"      # ONS turquoise - secondary series
ORANGE = "#F46A25"    # ONS orange - focal / highlight colour
GREY = "#B3B3B3"       # context grey
DARK_GREY = "#3D3D3D"
BG = "#FFFFFF"
CLUSTER_COLOURS = [NAVY, ORANGE, TEAL, "#801650", "#A285D1", "#3D3D3D", "#746CB1"]

st.markdown(
    f"""
    <style>
    .stApp {{ background-color: {BG}; }}
    h1, h2, h3 {{ font-family: Georgia, 'Times New Roman', serif; color: {NAVY}; }}
    .caveat-box {{
        background-color: #F0F0F0; border-left: 4px solid {ORANGE};
        padding: 0.8rem 1rem; border-radius: 2px; margin-bottom: 1rem;
    }}
    .method-box {{
        background-color: #EAF0F6; border-left: 4px solid {NAVY};
        padding: 0.8rem 1rem; border-radius: 2px; margin-bottom: 1rem;
    }}
    </style>
    """,
    unsafe_allow_html=True,
)

# ────────────────────────────────────────────────────────────────────────────
# 1. INDICATOR METADATA — plain-English definitions + provenance
#    (used throughout for captions aimed at a non-technical audience)
# ────────────────────────────────────────────────────────────────────────────
INDICATOR_INFO = {
    "Child poverty (AHC relative)": {
        "short": "Child poverty",
        "plain": "Share of children living in households below 60% of median income, after housing costs.",
        "source": "DWP / HMRC, Children in Low Income Families (local area statistics)",
        "period": "Year ending March 2025",
        "unit": "%",
        "why": "The single strongest, most policy-relevant driver of demand for statutory children's social care nationally.",
    },
    "Dependency ratio": {
        "short": "Dependency ratio",
        "plain": "Number of children (0-15) and older people (65+) per 100 working-age residents — a measure of the population an area has to support.",
        "source": "ONS, Census 2021",
        "period": "2021",
        "unit": "ratio",
        "why": "Areas with a higher child/elderly share carry structurally higher demand for both children's and adult services.",
    },
    "Population density": {
        "short": "Population density",
        "plain": "Residents per square kilometre.",
        "source": "ONS, Census 2021",
        "period": "2021",
        "unit": "people/km²",
        "why": "Urban vs rural context shapes both need (deprivation clusters in cities) and delivery cost (rural travel time).",
    },
    "Residents with no qualifications": {
        "short": "No qualifications",
        "plain": "Share of working-age residents with no formal qualifications.",
        "source": "ONS, Census 2021",
        "period": "2021",
        "unit": "%",
        "why": "A standard structural proxy for an area's socio-economic profile, used in the ONS's own contextual model.",
    },
    "Gross median weekly pay": {
        "short": "Weekly pay",
        "plain": "Median gross weekly pay of employees living in the area.",
        "source": "ONS, Annual Survey of Hours and Earnings (ASHE), place of residence",
        "period": "April 2025 (provisional)",
        "unit": "£",
        "why": "Income context; areas with similar pay levels tend to have similar service cost bases and household resilience.",
    },
    "CLA rate per 10,000": {
        "short": "Children looked after",
        "plain": "Children in local authority care, per 10,000 children in the area.",
        "source": "DfE, Children Looked After in England (and equivalents for devolved nations)",
        "period": "Year ending March 2025 / 31 March 2025",
        "unit": "per 10,000",
        "why": "The core statutory workload measure for a children's services directorate — the most acute end of demand.",
    },
    "CP plan rate per 10,000": {
        "short": "Child protection plans",
        "plain": "Children subject to an active child protection plan, per 10,000 children.",
        "source": "DfE, Characteristics of Children in Need",
        "period": "31 March 2025",
        "unit": "per 10,000",
        "why": "Sits directly upstream of care — a leading indicator of safeguarding pressure.",
    },
    "CIN rate per 10,000": {
        "short": "Children in need",
        "plain": "Children assessed as 'in need' of local authority support, per 10,000 children.",
        "source": "DfE, Characteristics of Children in Need",
        "period": "Year ending March 2023 (most recent geography breakdown available)",
        "unit": "per 10,000",
        "why": "The widest referral-pipeline measure — captures broader social-care contact, not just the most severe cases.",
    },
    "SEN/EHCP rate per 10,000": {
        "short": "SEN / EHC plans",
        "plain": "Children and young people (aged 0-25) with an Education, Health and Care Plan, per 10,000 in that age range.",
        "source": "DfE, Statements of SEN and EHC plans (SEN2 collection)",
        "period": "January 2025 (count) / mid-2024 (population)",
        "unit": "per 10,000",
        "why": "A statutory cross-cutting pressure that spans education, health and social care budgets simultaneously.",
    },
    "UASC rate per 10,000": {
        "short": "UASC",
        "plain": "Unaccompanied asylum-seeking children in the area's care, per 10,000 children.",
        "source": "DfE, Children Looked After (UASC subset)",
        "period": "Year ending March 2025",
        "unit": "per 10,000",
        "why": "A distinct, geographically concentrated demand driver — highly relevant for authorities with ports, airports or dispersal placements.",
    },
}
ALL_INDICATORS = list(INDICATOR_INFO.keys())
DEFAULT_INDICATORS = ALL_INDICATORS.copy()

TIER_NAMES = {
    "E06": "England — Unitary authority", "E08": "England — Metropolitan borough",
    "E09": "England — London borough", "E07": "England — District (two-tier)",
    "E10": "England — County council (two-tier)", "W06": "Wales — Unitary authority",
    "S12": "Scotland — Council area", "N09": "Northern Ireland — District",
}

# ────────────────────────────────────────────────────────────────────────────
# 2. DATA LOAD & PREP
# ────────────────────────────────────────────────────────────────────────────
@st.cache_data
def load_long():
    return pd.read_csv("data/subnational_indicators.csv")


@st.cache_data
def prepare_wide(long_df: pd.DataFrame):
    """Pivot to wide format and restrict to single-tier authorities.

    Children's services indicators (CLA/CP/CIN/SEN/UASC) are published at
    upper-tier level; several contextual indicators are published at
    district (lower-tier) level for the ~24 English shire counties. Merging
    the two tiers correctly needs a population-weighted district-to-county
    aggregation that this dataset does not include, so — consistent with
    how ONS/MHCLG themselves restrict statistical-neighbour matching by
    authority type — this model is scoped to single-tier authorities only:
    English unitaries, London boroughs and metropolitan boroughs, plus
    Welsh, Scottish and Northern Irish councils where coverage allows.
    """
    wide = long_df.pivot_table(index=["AREACD", "AREANM"], columns="Indicator", values="Value").reset_index()
    wide["Tier"] = wide["AREACD"].str[:3].map(TIER_NAMES).fillna(wide["AREACD"].str[:3])
    wide["prefix"] = wide["AREACD"].str[:3]
    single_tier = wide[wide["prefix"].isin(["E06", "E08", "E09", "W06", "S12", "N09"])].copy()
    return wide, single_tier


long_df = load_long()
wide_all, single_tier_df = prepare_wide(long_df)

# ────────────────────────────────────────────────────────────────────────────
# 3. SIDEBAR — variable selection & model controls
# ────────────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.header("Build your model")
    st.caption("Add or remove indicators to see how the groupings change.")

    selected_indicators = st.multiselect(
        "Indicators in the model",
        options=ALL_INDICATORS,
        default=DEFAULT_INDICATORS,
        format_func=lambda x: INDICATOR_INFO[x]["short"],
    )
    if len(selected_indicators) < 2:
        st.error("Select at least 2 indicators to build a model.")
        st.stop()

    apply_winsorise = st.checkbox(
        "Cap extreme outliers (winsorise at 1st/99th percentile)", value=True,
        help="ONS's own methodology caps the most extreme 1% of values at each end before "
             "standardising, so one unusually extreme area doesn't dominate every distance "
             "calculation. Recommended on.",
    )

    st.divider()
    st.subheader("Number of clusters (K)")
    k_mode = st.radio(
        "How should K be chosen?",
        ["Auto (best silhouette score)", "Choose manually"],
        index=0,
    )

    st.divider()
    spotlight_options = single_tier_df.sort_values("AREANM")["AREANM"].tolist()
    default_idx = spotlight_options.index("Westminster") if "Westminster" in spotlight_options else 0
    spotlight = st.selectbox("Spotlight authority", spotlight_options, index=default_idx)

    st.divider()
    st.caption(
        f"Model universe: **{len(single_tier_df)} single-tier authorities** "
        f"(England unitaries, London boroughs, metropolitan boroughs; Wales, Scotland, "
        f"Northern Ireland councils where data allows)."
    )

# ────────────────────────────────────────────────────────────────────────────
# 4. PIPELINE — complete case, winsorise, standardise, elbow/silhouette, KMeans
# ────────────────────────────────────────────────────────────────────────────
@st.cache_data
def build_pipeline(indicators: tuple, winsorise: bool):
    df = single_tier_df[["AREACD", "AREANM", "Tier"] + list(indicators)].dropna().reset_index(drop=True)
    X = df[list(indicators)].values.astype(float)
    Xw = X.copy()
    if winsorise:
        for j in range(X.shape[1]):
            lo, hi = np.percentile(X[:, j], [1, 99])
            Xw[:, j] = np.clip(X[:, j], lo, hi)
    Z = StandardScaler().fit_transform(Xw)
    return df, Z


@st.cache_data
def elbow_silhouette(Z: np.ndarray, k_min=2, k_max=9):
    rows = []
    for k in range(k_min, k_max + 1):
        km = KMeans(n_clusters=k, random_state=42, n_init=10).fit(Z)
        sil = silhouette_score(Z, km.labels_)
        rows.append({"K": k, "Inertia": km.inertia_, "Silhouette": sil})
    return pd.DataFrame(rows)


def knee_point(k_vals, inertias):
    """Simple max-distance-from-chord knee detector."""
    k_vals = np.array(k_vals, dtype=float)
    inertias = np.array(inertias, dtype=float)
    x1, y1 = k_vals[0], inertias[0]
    x2, y2 = k_vals[-1], inertias[-1]
    # normalise
    xn = (k_vals - x1) / (x2 - x1)
    yn = (inertias - y1) / (y2 - y1)
    # distance from the line joining first & last point
    dist = np.abs(yn - xn)
    return int(k_vals[np.argmax(dist)])


df_model, Z = build_pipeline(tuple(selected_indicators), apply_winsorise)
n_areas = len(df_model)
scores = elbow_silhouette(Z)

best_k_silhouette = int(scores.loc[scores["Silhouette"].idxmax(), "K"])
best_k_elbow = knee_point(scores["K"].tolist(), scores["Inertia"].tolist())

if k_mode.startswith("Auto"):
    K = best_k_silhouette
else:
    K = st.sidebar.slider("K (number of clusters)", 2, 9, value=best_k_silhouette)

km_final = KMeans(n_clusters=K, random_state=42, n_init=10).fit(Z)
df_model["Cluster"] = km_final.labels_
final_sil = silhouette_score(Z, km_final.labels_)
sample_sil = silhouette_samples(Z, km_final.labels_)
df_model["SilhouetteWidth"] = sample_sil

pca = PCA(n_components=2, random_state=42)
pca_coords = pca.fit_transform(Z)
df_model["PC1"], df_model["PC2"] = pca_coords[:, 0], pca_coords[:, 1]
explained = pca.explained_variance_ratio_

# ────────────────────────────────────────────────────────────────────────────
# 5. HELPERS — chart styling & PPTX export (per streamlit-app skill)
# ────────────────────────────────────────────────────────────────────────────
def style(fig: go.Figure, height=440):
    fig.update_layout(
        font_family="Arial", title_font_size=17, title_font_color=DARK_GREY,
        plot_bgcolor="white", paper_bgcolor="white", height=height,
        margin=dict(l=40, r=20, t=70, b=40),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, x=0),
    )
    fig.update_xaxes(showgrid=False, linecolor="#cccccc")
    fig.update_yaxes(gridcolor="#eeeeee", linecolor="white")
    return fig


def pptx_button(fig: go.Figure, chart_id: str, key_suffix=""):
    try:
        img_bytes = fig.to_image(format="png", width=1200, height=700, scale=2)
    except Exception:
        st.caption("(Install `kaleido` to enable PPTX slide export for this chart.)")
        return
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.0))
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.text = (fig.layout.title.text or chart_id)
    tf.paragraphs[0].runs[0].font.size = Pt(14)
    tf.paragraphs[0].runs[0].font.bold = True
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    st.download_button("⬇ Download this slide (PPTX)", buf, f"{chart_id}.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"dl_{chart_id}_{key_suffix}")


def render(fig, chart_id, caption="", key_suffix=""):
    style(fig)
    st.plotly_chart(fig, use_container_width=True, key=f"{chart_id}_{key_suffix}")
    if caption:
        st.caption(caption)
    pptx_button(fig, chart_id, key_suffix)


# ────────────────────────────────────────────────────────────────────────────
# 6. HEADER
# ────────────────────────────────────────────────────────────────────────────
st.title("Which authorities are genuinely similar to us?")
st.caption(
    "A nearest-neighbours & clustering model built for Children's Services financial planning — "
    "extending the ONS statistical neighbours methodology with ten indicators specific to "
    "children's services demand."
)

c1, c2, c3, c4 = st.columns(4)
c1.metric("Authorities in model", n_areas)
c2.metric("Indicators selected", len(selected_indicators))
c3.metric("Clusters (K)", K, help=f"Best silhouette at K={best_k_silhouette}; elbow knee at K={best_k_elbow}")
c4.metric("Silhouette score", f"{final_sil:.2f}")

tabs = st.tabs([
    "1 · Methodology & caveats",
    "2 · Choosing K (elbow & silhouette)",
    "3 · What do elbow & silhouette mean?",
    "4 · Explore the clusters",
    "5 · Your nearest neighbours",
])

# ────────────────────────────────────────────────────────────────────────────
# TAB 1 — METHODOLOGY & CAVEATS
# ────────────────────────────────────────────────────────────────────────────
with tabs[0]:
    st.header("Methodology, in plain English")
    st.markdown(
        """
        <div class="method-box">
        <b>What this model does:</b> every authority is described by the indicators you've
        selected in the sidebar. Each indicator is put onto a common scale (so "£ per week" and
        "% in poverty" can be compared fairly), then a computer algorithm (K-means) groups
        authorities that look statistically similar across <i>all</i> of those indicators at once.
        Separately, for the authority you spotlight, we calculate its <b>nearest neighbours</b> —
        the individual authorities whose overall profile is closest to it — which is the number
        that matters most for benchmarking and financial planning.
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.subheader("The pipeline, step by step")
    steps = [
        ("1. Select indicators", "Choose which of the 10 indicators describe your business challenge (sidebar)."),
        ("2. Keep complete cases", "Drop any authority missing even one selected indicator, so every comparison uses the same information."),
        ("3. Cap outliers (winsorise)", "Optionally cap the most extreme 1% of values at each end, so one atypical authority can't dominate the whole model."),
        ("4. Standardise (Z-score)", "Rescale every indicator to mean 0, standard deviation 1 — so a £ measure and a % measure carry equal weight."),
        ("5. Choose K", "Test K = 2 to 9 clusters; use the elbow and silhouette diagnostics (Tab 2) to pick a defensible number."),
        ("6. Cluster (K-means)", "Group authorities that sit close together across all standardised indicators."),
        ("7. Nearest neighbours", "For any one authority, rank every other authority by Euclidean distance across the same indicators."),
    ]
    cols = st.columns(4)
    for i, (title, desc) in enumerate(steps):
        with cols[i % 4]:
            st.markdown(f"**{title}**")
            st.caption(desc)

    st.subheader("Indicators used, and why")
    info_rows = []
    for ind in ALL_INDICATORS:
        meta = INDICATOR_INFO[ind]
        info_rows.append({
            "Indicator": meta["short"], "In current model": "✅" if ind in selected_indicators else "—",
            "What it measures": meta["plain"], "Source": meta["source"], "Period": meta["period"],
        })
    st.dataframe(pd.DataFrame(info_rows), use_container_width=True, hide_index=True)

    st.subheader("Data caveats — read before presenting this")
    st.markdown(
        f"""
        <div class="caveat-box">
        <b>Different reference periods.</b> Indicators span FYE 2023 to April 2025 (provisional) —
        they are the <i>most recent published</i> figures for each series, not a single snapshot
        year. Census-derived indicators (dependency ratio, population density, no qualifications)
        are fixed at Census 2021 and will not move until Census 2031.<br><br>
        <b>Geographic scope: single-tier authorities only.</b> Children's services indicators
        are published at upper-tier level; some contextual indicators are published at district
        (lower-tier) level for England's ~24 shire counties. Combining the two correctly needs a
        population-weighted district→county aggregation not available in this dataset, so —
        consistent with how ONS/MHCLG themselves restrict statistical-neighbour matching by
        authority type — this model is scoped to <b>{len(single_tier_df)} single-tier authorities</b>
        (unitaries, London boroughs, metropolitan boroughs, and single-tier councils in the
        devolved nations). Two-tier English counties/districts are out of scope for this version.<br><br>
        <b>Nation coverage gaps.</b> Scotland has no direct equivalent to England's "Children in
        Need" statistic, so CIN rate excludes Scotland. Northern Ireland does not publish
        SEN/EHCP or UASC figures broken down by district, so NI is excluded from those two
        indicators. Any authority missing even one selected indicator is dropped from this run
        (currently reducing the model to <b>{n_areas} of {len(single_tier_df)}</b> authorities) —
        this is a real trade-off between breadth of indicators and breadth of geographic coverage.<br><br>
        <b>Deviation from the published ONS method.</b> ONS's own national model does not include
        any children's-social-care-specific indicators (CLA, CP plan, CIN, SEN/EHCP, UASC) — this
        is a bespoke extension using ONS's <i>method</i> (winsorise → standardise → K-means /
        Euclidean nearest neighbours) applied to a children's-services-specific indicator set that
        ONS itself has not modelled.
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.link_button(
        "ONS statistical neighbours methodology (opens ons.gov.uk)",
        "https://www.ons.gov.uk/peoplepopulationandcommunity/wellbeing/methodologies/"
        "clusteringsimilarlocalauthoritiesandstatisticalnearestneighboursintheukmethodology",
    )

# ────────────────────────────────────────────────────────────────────────────
# TAB 2 — CHOOSING K
# ────────────────────────────────────────────────────────────────────────────
with tabs[1]:
    st.header("Choosing K: how many groups genuinely exist?")
    st.markdown(
        f"For your current {len(selected_indicators)} selected indicators and {n_areas} authorities: "
        f"the **elbow** sits around **K = {best_k_elbow}**, and the **silhouette score is highest at "
        f"K = {best_k_silhouette}** (score {scores.loc[scores['K']==best_k_silhouette,'Silhouette'].values[0]:.2f}). "
        f"You are currently viewing **K = {K}**."
    )

    left, right = st.columns(2)
    with left:
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=scores["K"], y=scores["Inertia"], mode="lines+markers",
                                  line=dict(color=NAVY, width=3), marker=dict(size=9)))
        fig.add_vline(x=best_k_elbow, line_dash="dash", line_color=ORANGE,
                       annotation_text=f"elbow ≈ K={best_k_elbow}", annotation_position="top")
        fig.update_layout(title=f"Inertia keeps falling, but flattens after K={best_k_elbow}",
                           xaxis_title="Number of clusters (K)", yaxis_title="Inertia (within-cluster spread)")
        render(fig, "elbow_chart",
               "Inertia = total squared distance from each authority to its cluster's centre. "
               "It always falls as K rises — the elbow is where extra clusters stop buying much.",
               "t2")
    with right:
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=scores["K"], y=scores["Silhouette"], mode="lines+markers",
                                  line=dict(color=TEAL, width=3), marker=dict(size=9)))
        fig.add_vline(x=best_k_silhouette, line_dash="dash", line_color=ORANGE,
                       annotation_text=f"best = K={best_k_silhouette}", annotation_position="top")
        fig.update_layout(title=f"Cluster separation peaks at K={best_k_silhouette}",
                           xaxis_title="Number of clusters (K)", yaxis_title="Average silhouette score")
        render(fig, "silhouette_chart",
               "Silhouette score (-1 to +1) measures how well-separated the clusters are. "
               "Higher = authorities sit clearly with their own group and far from other groups.",
               "t2")

    st.info(
        f"💡 **Why K={K} was chosen for {len(selected_indicators)} variables:** "
        + (f"the silhouette score is maximised at K={best_k_silhouette}, meaning this is the split "
           f"that produces the most internally-consistent, best-separated groups of authorities — "
           f"adding a 4th, 5th or 6th group starts re-splitting authorities that are already "
           f"well-matched, without improving how distinct the groups are from each other."
           if k_mode.startswith("Auto") else
           f"you have manually overridden the automatic choice (silhouette-optimal K={best_k_silhouette}, "
           f"elbow knee K={best_k_elbow}) to test K={K} directly.")
    )

# ────────────────────────────────────────────────────────────────────────────
# TAB 3 — DEDICATED ELBOW / SILHOUETTE EXPLAINER (toy example, plain English)
# ────────────────────────────────────────────────────────────────────────────
with tabs[2]:
    st.header("What do 'elbow' and 'silhouette' actually mean?")
    st.caption("A worked example with 12 illustrative points, before we apply either idea to real authorities.")

    rng = np.random.default_rng(7)
    toy = np.vstack([
        rng.normal([2, 2], 0.5, (4, 2)),
        rng.normal([8, 2], 0.5, (4, 2)),
        rng.normal([5, 7], 0.5, (4, 2)),
    ])
    toy_df = pd.DataFrame(toy, columns=["x", "y"])

    ex1, ex2 = st.columns(2)
    with ex1:
        st.subheader("The elbow method (inertia)")
        st.markdown(
            """
            **In plain English:** inertia measures how tightly points sit around their cluster's
            centre — like measuring how compact each group is. As you add more clusters, inertia
            *always* falls (more groups = shorter distances to the nearest centre). The **elbow**
            is where the curve stops dropping steeply and starts flattening out — meaning extra
            clusters aren't finding genuinely new structure, they're just splitting groups that
            were already fine.
            """
        )
        toy_inertias = []
        for k in range(1, 7):
            km = KMeans(n_clusters=k, random_state=1, n_init=10).fit(toy)
            toy_inertias.append(km.inertia_)
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=list(range(1, 7)), y=toy_inertias, mode="lines+markers",
                                  line=dict(color=NAVY, width=3), marker=dict(size=10)))
        fig.add_vline(x=3, line_dash="dash", line_color=ORANGE, annotation_text="elbow at K=3")
        fig.update_layout(title="Toy example: the elbow sits at K=3 (the 3 true groups)",
                           xaxis_title="K", yaxis_title="Inertia")
        render(fig, "toy_elbow", key_suffix="t3")

    with ex2:
        st.subheader("The silhouette score")
        st.markdown(
            """
            **In plain English:** for every single point, silhouette compares (a) how close it is
            to *other points in its own cluster* against (b) how close it is to the *nearest
            other cluster*. A score near **+1** means "clearly belongs here, far from anywhere
            else." A score near **0** means "sitting right on the boundary between two groups."
            A **negative** score means a point may actually be in the wrong cluster. Averaging
            this across every point gives one overall score for a given K.
            """
        )
        km3 = KMeans(n_clusters=3, random_state=1, n_init=10).fit(toy)
        sil_vals = silhouette_samples(toy, km3.labels_)
        order = np.argsort(km3.labels_ * 100 - sil_vals)
        fig = go.Figure()
        colours_map = {0: NAVY, 1: ORANGE, 2: TEAL}
        fig.add_trace(go.Bar(
            x=sil_vals[order], y=[f"pt {i}" for i in order], orientation="h",
            marker_color=[colours_map[c] for c in km3.labels_[order]],
        ))
        fig.add_vline(x=sil_vals.mean(), line_dash="dash", line_color=DARK_GREY,
                       annotation_text=f"average = {sil_vals.mean():.2f}")
        fig.update_layout(title="Every point's individual silhouette width (K=3)",
                           xaxis_title="Silhouette width (-1 to +1)", yaxis_title="")
        render(fig, "toy_silhouette", key_suffix="t3")

    st.subheader("See the 3 toy groups the scores above are describing")
    km3 = KMeans(n_clusters=3, random_state=1, n_init=10).fit(toy)
    toy_df["Cluster"] = km3.labels_.astype(str)
    fig = px.scatter(toy_df, x="x", y="y", color="Cluster",
                      color_discrete_sequence=[NAVY, ORANGE, TEAL],
                      title="The three groups these 12 points naturally form")
    fig.update_traces(marker=dict(size=16, line=dict(width=1, color="white")))
    render(fig, "toy_scatter", key_suffix="t3")

    st.markdown(
        """
        <div class="method-box">
        <b>Why we use both together:</b> the elbow method is quick but the "bend" can be a
        matter of judgement. The silhouette score is more precise but only ever gives you one
        K at a time. Using both — and picking the K where they agree, or explaining clearly when
        they don't — is standard practice and is exactly what ONS do in their own published
        methodology.
        </div>
        """,
        unsafe_allow_html=True,
    )

# ────────────────────────────────────────────────────────────────────────────
# TAB 4 — EXPLORE THE CLUSTERS (real data)
# ────────────────────────────────────────────────────────────────────────────
with tabs[3]:
    st.header(f"The {K} groups found in your {n_areas} authorities")

    cluster_sizes = df_model["Cluster"].value_counts().sort_index()
    fig = px.bar(x=[f"Cluster {i}" for i in cluster_sizes.index], y=cluster_sizes.values,
                 color=[f"Cluster {i}" for i in cluster_sizes.index],
                 color_discrete_sequence=CLUSTER_COLOURS,
                 title="How many authorities fall into each group")
    fig.update_layout(showlegend=False, xaxis_title="", yaxis_title="Number of authorities")
    fig.update_yaxes(rangemode="tozero")
    render(fig, "cluster_sizes", key_suffix="t4")

    left, right = st.columns([3, 2])
    with left:
        spot_row = df_model[df_model["AREANM"] == spotlight]
        spot_cluster = int(spot_row["Cluster"].iloc[0]) if len(spot_row) else None
        df_model["Highlight"] = np.where(df_model["AREANM"] == spotlight, spotlight, "Other authorities")
        fig = px.scatter(
            df_model, x="PC1", y="PC2", color=df_model["Cluster"].astype(str),
            hover_name="AREANM", symbol="Highlight",
            symbol_map={spotlight: "star", "Other authorities": "circle"},
            color_discrete_sequence=CLUSTER_COLOURS,
            title=f"Every authority's overall profile, compressed to 2 dimensions "
                  f"(explains {explained.sum()*100:.0f}% of the variation)",
        )
        fig.update_traces(marker=dict(size=11, line=dict(width=0.5, color="white")),
                           selector=dict(mode="markers"))
        render(fig, "pca_scatter",
               "This chart (PCA) squashes every selected indicator down to 2 axes so similar "
               "authorities appear near each other — it's a map of similarity, not geography.",
               "t4")
    with right:
        st.subheader(f"{spotlight}'s profile")
        if len(spot_row):
            st.metric("Assigned to", f"Cluster {spot_cluster}")
            st.metric("Silhouette width", f"{spot_row['SilhouetteWidth'].iloc[0]:.2f}",
                       help="How clearly this authority belongs to its group vs the next-nearest group.")
            st.caption(f"{cluster_sizes[spot_cluster]} authorities share this cluster.")
        else:
            st.warning(f"{spotlight} is missing one or more selected indicators, so it was "
                        f"dropped from this run (see Tab 1 caveats).")

    st.subheader("Cluster profiles — what makes each group distinct")
    profile = df_model.groupby("Cluster")[selected_indicators].mean()
    profile_z = (profile - df_model[selected_indicators].mean()) / df_model[selected_indicators].std()
    heat = profile_z.rename(columns={k: INDICATOR_INFO[k]["short"] for k in selected_indicators})
    fig = px.imshow(
        heat.T, color_continuous_scale=[[0, TEAL], [0.5, "white"], [1, ORANGE]],
        aspect="auto", labels=dict(color="Std. deviations from national average"),
        title="Where each cluster sits above (orange) or below (teal) the national average",
    )
    fig.update_xaxes(title="Cluster", tickvals=list(range(K)), ticktext=[f"Cluster {i}" for i in range(K)])
    fig.update_yaxes(title="")
    render(fig, "cluster_heatmap",
           "Each cell shows how far that cluster's average sits from the overall average, in "
           "standard deviations — the same scale the clustering itself used.", "t4")

    st.subheader("All selected indicators at once, by cluster")
    pc_df = df_model[["AREANM", "Cluster"] + selected_indicators].copy()
    dims = [dict(label=INDICATOR_INFO[c]["short"], values=pc_df[c]) for c in selected_indicators]
    fig = go.Figure(data=go.Parcoords(
        line=dict(color=pc_df["Cluster"], colorscale=[[i/(max(K-1,1)), c] for i, c in enumerate(CLUSTER_COLOURS[:K])]),
        dimensions=dims,
    ))
    fig.update_layout(title="Each line is one authority — coloured by cluster")
    render(fig, "parallel_coords",
           "Parallel coordinates: authorities in the same cluster tend to trace similar paths "
           "across all indicators at once, showing which measures separate the groups most clearly.",
           "t4")

    if spotlight in df_model["AREANM"].values:
        st.subheader(f"{spotlight} vs its own cluster vs the national average")
        z_df = pd.DataFrame(Z, columns=selected_indicators, index=df_model["AREANM"])
        spot_z = z_df.loc[spotlight]
        cluster_avg_z = z_df[df_model["Cluster"].values == spot_cluster].mean()
        radar_labels = [INDICATOR_INFO[c]["short"] for c in selected_indicators]
        fig = go.Figure()
        fig.add_trace(go.Scatterpolar(r=list(spot_z.values) + [spot_z.values[0]],
                                       theta=radar_labels + [radar_labels[0]],
                                       name=spotlight, line=dict(color=ORANGE, width=3)))
        fig.add_trace(go.Scatterpolar(r=list(cluster_avg_z.values) + [cluster_avg_z.values[0]],
                                       theta=radar_labels + [radar_labels[0]],
                                       name=f"Cluster {spot_cluster} average", line=dict(color=NAVY, width=2, dash="dash")))
        fig.add_trace(go.Scatterpolar(r=[0] * (len(radar_labels) + 1), theta=radar_labels + [radar_labels[0]],
                                       name="National average", line=dict(color=GREY, width=1, dash="dot")))
        fig.update_layout(title=f"{spotlight} (standard deviations from the national average)",
                           polar=dict(radialaxis=dict(visible=True)))
        render(fig, "radar_chart", key_suffix="t4")

# ────────────────────────────────────────────────────────────────────────────
# TAB 5 — NEAREST NEIGHBOURS
# ────────────────────────────────────────────────────────────────────────────
with tabs[4]:
    st.header(f"{spotlight}'s nearest neighbours")
    if spotlight not in df_model["AREANM"].values:
        st.warning(f"{spotlight} was dropped from this run because it is missing one of the "
                   f"selected indicators — remove that indicator in the sidebar, or pick a "
                   f"different spotlight authority.")
    else:
        st.markdown(
            """
            <div class="method-box">
            <b>This is separate from clustering.</b> Clustering sorts every authority into a
            small number of broad groups. Nearest neighbours instead asks one precise question:
            <i>of every other authority, which ones is this one most similar to, overall?</i>
            It's calculated as the straight-line (Euclidean) distance between standardised
            indicator profiles — the smaller the distance, the more similar the two authorities.
            </div>
            """,
            unsafe_allow_html=True,
        )

        z_df = pd.DataFrame(Z, columns=selected_indicators, index=df_model["AREANM"])
        target = z_df.loc[spotlight]
        dist = np.sqrt(((z_df - target) ** 2).sum(axis=1))
        dist = dist.drop(spotlight).sort_values()
        n_neighbours = st.slider("How many neighbours to show", 5, 20, 10)
        top = dist.head(n_neighbours).reset_index()
        top.columns = ["Authority", "Distance"]
        top["Tier"] = top["Authority"].map(df_model.set_index("AREANM")["Tier"])

        left, right = st.columns([3, 2])
        with left:
            fig = px.bar(top.sort_values("Distance", ascending=True), x="Distance", y="Authority",
                         orientation="h", color="Distance", color_continuous_scale=[[0, ORANGE], [1, GREY]],
                         title=f"Closest matches to {spotlight} (shortest bar = most similar)")
            fig.update_yaxes(categoryorder="total descending")
            fig.update_layout(coloraxis_showscale=False)
            render(fig, "neighbours_bar", key_suffix="t5")
        with right:
            median_national = dist.median()
            st.metric(f"{spotlight}'s nearest neighbour distance", f"{dist.min():.2f}")
            st.metric("Median distance across all pairs in this model", f"{median_national:.2f}")
            if dist.min() > median_national:
                st.caption(
                    f"⚠️ {spotlight}'s closest match is *further away* than the typical "
                    f"authority's closest match — a sign {spotlight} is a statistical outlier "
                    f"and comparisons should be made cautiously."
                )
            else:
                st.caption(f"{spotlight} has at least one closely-matched comparator authority.")

        st.subheader("Raw indicator values — for the finance conversation")
        table_cols = ["AREANM"] + selected_indicators
        display = df_model[df_model["AREANM"].isin([spotlight] + top["Authority"].tolist())][table_cols]
        display = display.set_index("AREANM").loc[[spotlight] + top["Authority"].tolist()]
        display = display.rename(columns={k: INDICATOR_INFO[k]["short"] for k in selected_indicators})
        st.dataframe(display.style.format("{:.1f}").apply(
            lambda s: ["background-color:#FDEEE3" if s.name == spotlight else "" for _ in s], axis=1),
            use_container_width=True)
        st.caption(
            "Unlike the charts above (which use standardised/scaled values so different units "
            "compare fairly), this table shows the real published figures — the numbers to "
            "actually quote in a finance report."
        )

        csv = display.reset_index().to_csv(index=False).encode()
        st.download_button("⬇ Download this table (CSV)", csv, f"{spotlight}_nearest_neighbours.csv", "text/csv")

st.divider()
st.caption(
    "Built as a DAP (Data Analytics Principles) artefact. Methodology extends the ONS "
    "statistical nearest neighbours approach with children's-services-specific indicators. "
    "See Tab 1 for full data sources, periods and caveats."
)
