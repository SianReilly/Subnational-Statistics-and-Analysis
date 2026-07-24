# pip install streamlit plotly pandas scikit-learn python-pptx kaleido openpyxl numpy
"""
Children's Services Nearest Neighbours & Clustering Explorer
--------------------------------------------------------------
A DAP (Data Analytics Principles) artefact built for a Local Authority
Children's Services directorate. Extends the ONS statistical nearest
neighbours / clustering methodology
(https://www.ons.gov.uk/peoplepopulationandcommunity/wellbeing/methodologies/
clusteringsimilarlocalauthoritiesandstatisticalnearestneighboursintheukmethodology)
with ten indicators specific to children's services demand.

Run with:  streamlit run app.py
Data expected at:  data/subnational_indicators.csv  (long format:
AREACD, AREANM, Indicator, Period, Measure, Unit, Value)
"""

import io
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
from sklearn.metrics import silhouette_samples, silhouette_score
from sklearn.preprocessing import StandardScaler

# ────────────────────────────────────────────────────────────────────────────
# 0. PAGE CONFIG & STYLE (ONS / Economist palette)
# ────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Children's Services Nearest Neighbours",
    page_icon="🧭",
    layout="wide",
    initial_sidebar_state="expanded",
)

NAVY = "#12436D"      # ONS dark blue - primary
TEAL = "#28A197"      # ONS turquoise - secondary series
ORANGE = "#F46A25"    # ONS orange - focal / highlight colour
GREY = "#B3B3B3"       # context grey
DARK_GREY = "#3D3D3D"
BG = "#FFFFFF"
CLUSTER_COLOURS = [NAVY, ORANGE, TEAL, "#801650", "#A285D1", "#3D3D3D", "#746CB1"]

st.markdown(
    f"""
    <style>
    .stApp {{ background-color: {BG}; }}
    h1, h2, h3 {{ font-family: Georgia, 'Times New Roman', serif; color: {NAVY}; }}
    .caveat-box {{
        background-color: #F0F0F0; border-left: 4px solid {ORANGE};
        padding: 0.8rem 1rem; border-radius: 2px; margin-bottom: 1rem;
    }}
    .method-box {{
        background-color: #EAF0F6; border-left: 4px solid {NAVY};
        padding: 0.8rem 1rem; border-radius: 2px; margin-bottom: 1rem;
    }}
    </style>
    """,
    unsafe_allow_html=True,
)

# ────────────────────────────────────────────────────────────────────────────
# 1. INDICATOR METADATA — plain-English definitions + provenance
#    (used throughout for captions aimed at a non-technical audience)
# ────────────────────────────────────────────────────────────────────────────
INDICATOR_INFO = {
    "Child poverty (AHC relative)": {
        "short": "Child poverty",
        "plain": "Share of children living in households below 60% of median income, after housing costs.",
        "source": "DWP / HMRC, Children in Low Income Families (local area statistics)",
        "period": "Year ending March 2025",
        "unit": "%",
        "why": "The single strongest, most policy-relevant driver of demand for statutory children's social care nationally.",
    },
    "Dependency ratio": {
        "short": "Dependency ratio",
        "plain": "Number of children (0-15) and older people (65+) per 100 working-age residents — a measure of the population an area has to support.",
        "source": "ONS, Census 2021",
        "period": "2021",
        "unit": "ratio",
        "why": "Areas with a higher child/elderly share carry structurally higher demand for both children's and adult services.",
    },
    "Population density": {
        "short": "Population density",
        "plain": "Residents per square kilometre.",
        "source": "ONS, Census 2021",
        "period": "2021",
        "unit": "people/km²",
        "why": "Urban vs rural context shapes both need (deprivation clusters in cities) and delivery cost (rural travel time).",
    },
    "Residents with no qualifications": {
        "short": "No qualifications",
        "plain": "Share of working-age residents with no formal qualifications.",
        "source": "ONS, Census 2021",
        "period": "2021",
        "unit": "%",
        "why": "A standard structural proxy for an area's socio-economic profile, used in the ONS's own contextual model.",
    },
    "Gross median weekly pay": {
        "short": "Weekly pay",
        "plain": "Median gross weekly pay of employees living in the area.",
        "source": "ONS, Annual Survey of Hours and Earnings (ASHE), place of residence",
        "period": "April 2025 (provisional)",
        "unit": "£",
        "why": "Income context; areas with similar pay levels tend to have similar service cost bases and household resilience.",
    },
    "CLA rate per 10,000": {
        "short": "Children looked after",
        "plain": "Children in local authority care, per 10,000 children in the area.",
        "source": "DfE, Children Looked After in England (and equivalents for devolved nations)",
        "period": "Year ending March 2025 / 31 March 2025",
        "unit": "per 10,000",
        "why": "The core statutory workload measure for a children's services directorate — the most acute end of demand.",
    },
    "CP plan rate per 10,000": {
        "short": "Child protection plans",
        "plain": "Children subject to an active child protection plan, per 10,000 children.",
        "source": "DfE, Characteristics of Children in Need",
        "period": "31 March 2025",
        "unit": "per 10,000",
        "why": "Sits directly upstream of care — a leading indicator of safeguarding pressure.",
    },
    "CIN rate per 10,000": {
        "short": "Children in need",
        "plain": "Children assessed as 'in need' of local authority support, per 10,000 children.",
        "source": "DfE, Characteristics of Children in Need",
        "period": "Year ending March 2023 (most recent geography breakdown available)",
        "unit": "per 10,000",
        "why": "The widest referral-pipeline measure — captures broader social-care contact, not just the most severe cases.",
    },
    "SEN/EHCP rate per 10,000": {
        "short": "SEN / EHC plans",
        "plain": "Children and young people (aged 0-25) with an Education, Health and Care Plan, per 10,000 in that age range.",
        "source": "DfE, Statements of SEN and EHC plans (SEN2 collection)",
        "period": "January 2025 (count) / mid-2024 (population)",
        "unit": "per 10,000",
        "why": "A statutory cross-cutting pressure that spans education, health and social care budgets simultaneously.",
    },
    "UASC rate per 10,000": {
        "short": "UASC",
        "plain": "Unaccompanied asylum-seeking children in the area's care, per 10,000 children.",
        "source": "DfE, Children Looked After (UASC subset)",
        "period": "Year ending March 2025",
        "unit": "per 10,000",
        "why": "A distinct, geographically concentrated demand driver — highly relevant for authorities with ports, airports or dispersal placements.",
    },
}
ALL_INDICATORS = list(INDICATOR_INFO.keys())
DEFAULT_INDICATORS = ALL_INDICATORS.copy()

# Which of the 4 UK nations each indicator is actually published for (checked directly
# against the data — this is what drives the Scotland/NI exclusions in Tab 1 below).
NATION_COVERAGE = {
    "Child poverty (AHC relative)": {"England", "Wales", "Scotland", "N. Ireland"},
    "Dependency ratio": {"England", "Wales", "Scotland", "N. Ireland"},
    "Population density": {"England", "Wales", "Scotland", "N. Ireland"},
    "Residents with no qualifications": {"England", "Wales", "Scotland", "N. Ireland"},
    "Gross median weekly pay": {"England", "Wales", "Scotland", "N. Ireland"},
    "CLA rate per 10,000": {"England", "Wales", "Scotland", "N. Ireland"},
    "CP plan rate per 10,000": {"England", "Wales", "Scotland", "N. Ireland"},
    "CIN rate per 10,000": {"England", "Wales", "N. Ireland"},          # no Scottish equivalent
    "SEN/EHCP rate per 10,000": {"England", "Wales"},                    # no Scotland or NI breakdown
    "UASC rate per 10,000": {"England", "Wales", "Scotland"},            # no NI breakdown
}
ALL_NATIONS = {"England", "Wales", "Scotland", "N. Ireland"}

# Whether each indicator is part of ONS's own published national model, or a bespoke
# addition for this children's-services extension.
ONS_METHOD_DIFF = {
    "Child poverty (AHC relative)": "**Different measure basis** — this uses child poverty "
        "after housing costs (AHC); ONS's own published baseline model uses the before-housing-"
        "costs (BHC) measure. AHC is arguably more relevant for a Children's Services cost "
        "conversation (it reflects what a household actually has left to live on), but it is "
        "not a direct like-for-like swap with ONS's indicator of the same name.",
    "Dependency ratio": "Same as ONS's own baseline model — no change.",
    "Population density": "Same as ONS's own baseline model — no change.",
    "Residents with no qualifications": "Same as ONS's own baseline model — no change.",
    "Gross median weekly pay": "Same as ONS's own baseline model — no change.",
    "CLA rate per 10,000": "Bespoke addition — ONS's national model has no children's social care indicators.",
    "CP plan rate per 10,000": "Bespoke addition — ONS's national model has no children's social care indicators.",
    "CIN rate per 10,000": "Bespoke addition — ONS's national model has no children's social care indicators.",
    "SEN/EHCP rate per 10,000": "Bespoke addition — ONS's national model has no children's social care indicators.",
    "UASC rate per 10,000": "Bespoke addition — ONS's national model has no children's social care indicators.",
}

TIER_NAMES = {
    "E06": "England — Unitary authority", "E08": "England — Metropolitan borough",
    "E09": "England — London borough", "E07": "England — District (two-tier)",
    "E10": "England — County council (two-tier)", "W06": "Wales — Unitary authority",
    "S12": "Scotland — Council area", "N09": "Northern Ireland — District",
}

# ────────────────────────────────────────────────────────────────────────────
# 2. DATA LOAD & PREP
# ────────────────────────────────────────────────────────────────────────────
@st.cache_data
def load_long():
    import os
    # Works whether the CSV sits next to app.py (repo root) or in a data/ subfolder
    for candidate in ["subnational_indicators.csv", "data/subnational_indicators.csv"]:
        if os.path.exists(candidate):
            return pd.read_csv(candidate)
    raise FileNotFoundError(
        "Couldn't find subnational_indicators.csv next to app.py or in a data/ subfolder."
    )


@st.cache_data
def prepare_wide(long_df: pd.DataFrame):
    """Pivot to wide format and restrict to single-tier authorities.

    Children's services indicators (CLA/CP/CIN/SEN/UASC) are published at
    upper-tier level; several contextual indicators are published at
    district (lower-tier) level for the ~24 English shire counties. Merging
    the two tiers correctly needs a population-weighted district-to-county
    aggregation that this dataset does not include, so — consistent with
    how ONS/MHCLG themselves restrict statistical-neighbour matching by
    authority type — this model is scoped to single-tier authorities only:
    English unitaries, London boroughs and metropolitan boroughs, plus
    Welsh, Scottish and Northern Irish councils where coverage allows.
    """
    wide = long_df.pivot_table(index=["AREACD", "AREANM"], columns="Indicator", values="Value").reset_index()
    wide["Tier"] = wide["AREACD"].str[:3].map(TIER_NAMES).fillna(wide["AREACD"].str[:3])
    wide["prefix"] = wide["AREACD"].str[:3]
    single_tier = wide[wide["prefix"].isin(["E06", "E08", "E09", "W06", "S12", "N09"])].copy()
    return wide, single_tier


long_df = load_long()
wide_all, single_tier_df = prepare_wide(long_df)


@st.cache_data
def load_boundaries():
    """Loads an optional user-supplied GeoJSON boundary file for a real choropleth.
    Not required — the app works fully without it (see Tab 1)."""
    import json
    import os
    for candidate in ["data/la_boundaries.geojson", "la_boundaries.geojson"]:
        if os.path.exists(candidate):
            with open(candidate) as f:
                gj = json.load(f)
            props = gj.get("features", [{}])[0].get("properties", {})
            key_field = next((k for k in props if k.upper().startswith("LAD") and k.upper().endswith("CD")), None)
            if key_field is None:
                key_field = next((k for k in props if "CD" in k.upper()), None)
            gj["_key"] = f"properties.{key_field}" if key_field else "properties.LAD22CD"
            return gj
    return None

# ────────────────────────────────────────────────────────────────────────────
# 3. SIDEBAR — variable selection & model controls
# ────────────────────────────────────────────────────────────────────────────
CONTEXT_INDICATORS = [i for i in ALL_INDICATORS if not any(
    i.startswith(p) for p in ("CLA", "CP plan", "CIN", "SEN", "UASC"))]
CS_INDICATORS = [i for i in ALL_INDICATORS if i not in CONTEXT_INDICATORS]

with st.sidebar:
    st.markdown(
        f"""
        <div style="background-color:{NAVY}; padding:0.9rem 1rem; border-radius:4px; margin-bottom:1rem;">
        <span style="color:white; font-size:0.95rem; font-weight:bold;">
        🧭 Westminster City Council</span><br>
        <span style="color:#CFE0EC; font-size:0.85rem;">Children's Services Directorate — financial planning tool</span>
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.header("① Build your model")
    st.caption("Tick the indicators that should define 'similar' — every chart in the app updates live.")

    st.markdown("**Children's services indicators** *(the demand measures)*")
    cs_selected = st.multiselect(
        "Children's services indicators", options=CS_INDICATORS,
        default=CS_INDICATORS, format_func=lambda x: INDICATOR_INFO[x]["short"],
        label_visibility="collapsed",
        help="Direct measures of children's social care demand: care, protection plans, "
             "need, SEN/EHC plans, unaccompanied asylum-seeking children.",
    )
    st.markdown("**Context indicators** *(the area's wider profile)*")
    ctx_selected = st.multiselect(
        "Context indicators", options=CONTEXT_INDICATORS,
        default=CONTEXT_INDICATORS, format_func=lambda x: INDICATOR_INFO[x]["short"],
        label_visibility="collapsed",
        help="General socio-economic and demographic context, matching ONS's own baseline model.",
    )
    selected_indicators = cs_selected + ctx_selected
    st.caption(f"**{len(selected_indicators)} of {len(ALL_INDICATORS)} indicators selected.**")
    if len(selected_indicators) < 2:
        st.error("Select at least 2 indicators to build a model.")
        st.stop()

    apply_winsorise = st.checkbox(
        "Cap extreme outliers (winsorise at 1st/99th percentile)", value=True,
        help="ONS's own methodology caps the most extreme 1% of values at each end before "
             "standardising, so one unusually extreme area doesn't dominate every distance "
             "calculation. Recommended on.",
    )

    st.divider()
    st.header("② Number of clusters (K)")
    k_mode = st.radio(
        "How should K be chosen?",
        ["Auto (best silhouette score)", "Choose manually"],
        index=0,
        help="Auto picks the K with the highest silhouette score (see Tab 2). Choose manually "
             "to test a specific number of groups yourself.",
    )

    st.divider()
    st.header("③ Spotlight authority")
    spotlight_options = single_tier_df.sort_values("AREANM")["AREANM"].tolist()
    default_idx = spotlight_options.index("Westminster") if "Westminster" in spotlight_options else 0
    spotlight = st.selectbox(
        "Which authority to profile in Tabs 4 & 5", spotlight_options, index=default_idx,
        help="Every chart in Tabs 4 and 5 is written around this authority — change it to "
             "benchmark a different council.",
    )

    st.divider()
    st.caption(
        f"Model universe: **{len(single_tier_df)} single-tier authorities** "
        f"(England unitaries, London boroughs, metropolitan boroughs; Wales, Scotland, "
        f"Northern Ireland councils where data allows). See Tab 1 for full caveats."
    )

# ────────────────────────────────────────────────────────────────────────────
# 4. PIPELINE — complete case, winsorise, standardise, elbow/silhouette, KMeans
# ────────────────────────────────────────────────────────────────────────────
@st.cache_data
def build_pipeline(indicators: tuple, winsorise: bool):
    df = single_tier_df[["AREACD", "AREANM", "Tier"] + list(indicators)].dropna().reset_index(drop=True)
    X = df[list(indicators)].values.astype(float)
    Xw = X.copy()
    if winsorise:
        for j in range(X.shape[1]):
            lo, hi = np.percentile(X[:, j], [1, 99])
            Xw[:, j] = np.clip(X[:, j], lo, hi)
    Z = StandardScaler().fit_transform(Xw)
    return df, Z


@st.cache_data
def elbow_silhouette(Z: np.ndarray, k_min=2, k_max=9):
    rows = []
    for k in range(k_min, k_max + 1):
        km = KMeans(n_clusters=k, random_state=42, n_init=10).fit(Z)
        sil = silhouette_score(Z, km.labels_)
        rows.append({"K": k, "Inertia": km.inertia_, "Silhouette": sil})
    return pd.DataFrame(rows)


def knee_point(k_vals, inertias):
    """Simple max-distance-from-chord knee detector."""
    k_vals = np.array(k_vals, dtype=float)
    inertias = np.array(inertias, dtype=float)
    x1, y1 = k_vals[0], inertias[0]
    x2, y2 = k_vals[-1], inertias[-1]
    # normalise
    xn = (k_vals - x1) / (x2 - x1)
    yn = (inertias - y1) / (y2 - y1)
    # distance from the line joining first & last point
    dist = np.abs(yn - xn)
    return int(k_vals[np.argmax(dist)])


@st.cache_data
def other_validation_scores(Z: np.ndarray, k_min=2, k_max=9):
    """Calinski-Harabasz (higher=better) and Davies-Bouldin (lower=better) —
    two more internal validation indices, independent of silhouette's method,
    used to cross-check the choice of K."""
    from sklearn.metrics import calinski_harabasz_score, davies_bouldin_score
    rows = []
    for k in range(k_min, k_max + 1):
        km = KMeans(n_clusters=k, random_state=42, n_init=10).fit(Z)
        rows.append({
            "K": k,
            "Calinski-Harabasz": calinski_harabasz_score(Z, km.labels_),
            "Davies-Bouldin": davies_bouldin_score(Z, km.labels_),
        })
    return pd.DataFrame(rows)


@st.cache_data
def hopkins_statistic(Z: np.ndarray, sample_ratio=0.2, seed=42):
    """Clustering tendency check: is there genuine cluster structure in this data
    at all, before we even ask what K should be? H close to 1 = highly clusterable,
    H around 0.5 = essentially random/uniform (clustering would be spurious)."""
    from sklearn.neighbors import NearestNeighbors
    rng = np.random.default_rng(seed)
    n, d = Z.shape
    m = max(2, int(sample_ratio * n))
    idx = rng.choice(n, m, replace=False)
    Z_sample = Z[idx]
    mins, maxs = Z.min(axis=0), Z.max(axis=0)
    random_points = rng.uniform(mins, maxs, size=(m, d))
    nbrs = NearestNeighbors(n_neighbors=2).fit(Z)
    u_dist, _ = nbrs.kneighbors(random_points, n_neighbors=1)
    w_dist, _ = nbrs.kneighbors(Z_sample, n_neighbors=2)
    u_sum = u_dist[:, 0].sum()
    w_sum = w_dist[:, 1].sum()
    return u_sum / (u_sum + w_sum)


@st.cache_data
def bootstrap_stability(Z: np.ndarray, K: int, n_boot=30, frac=0.8, seed=42):
    """Repeatedly re-clusters random 80% subsamples and compares the result to the
    full-data clustering (restricted to the shared authorities) with the Adjusted
    Rand Index. Consistently high ARI = the clusters are a robust structural
    feature of the data, not an artefact of exactly which authorities happen to
    be included this run."""
    from sklearn.metrics import adjusted_rand_score
    rng = np.random.default_rng(seed)
    n = Z.shape[0]
    full_labels = KMeans(n_clusters=K, random_state=42, n_init=10).fit(Z).labels_
    aris = []
    for b in range(n_boot):
        idx = rng.choice(n, size=int(n * frac), replace=False)
        sub_labels = KMeans(n_clusters=K, random_state=seed + b, n_init=10).fit(Z[idx]).labels_
        aris.append(adjusted_rand_score(full_labels[idx], sub_labels))
    return np.array(aris)


df_model, Z = build_pipeline(tuple(selected_indicators), apply_winsorise)
n_areas = len(df_model)
scores = elbow_silhouette(Z)

best_k_silhouette = int(scores.loc[scores["Silhouette"].idxmax(), "K"])
best_k_elbow = knee_point(scores["K"].tolist(), scores["Inertia"].tolist())

if k_mode.startswith("Auto"):
    K = best_k_silhouette
else:
    K = st.sidebar.slider("K (number of clusters)", 2, 9, value=best_k_silhouette)

km_final = KMeans(n_clusters=K, random_state=42, n_init=10).fit(Z)
df_model["Cluster"] = km_final.labels_
final_sil = silhouette_score(Z, km_final.labels_)
sample_sil = silhouette_samples(Z, km_final.labels_)
df_model["SilhouetteWidth"] = sample_sil

pca = PCA(n_components=2, random_state=42)
pca_coords = pca.fit_transform(Z)
df_model["PC1"], df_model["PC2"] = pca_coords[:, 0], pca_coords[:, 1]
explained = pca.explained_variance_ratio_

# ────────────────────────────────────────────────────────────────────────────
# 5. HELPERS — chart styling & PPTX export (per streamlit-app skill)
# ────────────────────────────────────────────────────────────────────────────
def style(fig: go.Figure, height=460, legend_below=True):
    """Consistent, overlap-free styling: title top-left, legend pushed clear
    below the plot (never over the title or the data), generous auto-margins."""
    fig.update_layout(
        font_family="Arial", font_size=13,
        title=dict(font=dict(size=16, color=DARK_GREY), x=0.0, xanchor="left",
                    y=0.97, yanchor="top", pad=dict(b=10)),
        plot_bgcolor="white", paper_bgcolor="white", height=height,
        margin=dict(l=70, r=40, t=70, b=90 if legend_below else 50),
        legend=(dict(orientation="h", yanchor="top", y=-0.22, x=0, xanchor="left",
                      font=dict(size=11))
                if legend_below else dict(font=dict(size=11))),
        uniformtext_minsize=10,
    )
    fig.update_xaxes(showgrid=False, linecolor="#cccccc", automargin=True,
                      title_font=dict(size=12), tickfont=dict(size=11))
    fig.update_yaxes(gridcolor="#eeeeee", linecolor="white", automargin=True,
                      title_font=dict(size=12), tickfont=dict(size=11))
    return fig


def annotate_vline(fig, x, text, color=ORANGE, y_paper=1.0):
    """Adds a vertical line whose label sits in a clear band above the plot
    area (paper coordinates) so it never collides with the data or the title."""
    fig.add_vline(x=x, line_dash="dash", line_color=color, line_width=2)
    fig.add_annotation(
        x=x, y=y_paper, xref="x", yref="paper", text=text, showarrow=False,
        font=dict(size=12, color=color), bgcolor="white", bordercolor=color,
        borderwidth=1, borderpad=3, yanchor="bottom",
    )
    return fig


def pptx_button(fig: go.Figure, chart_id: str, key_suffix=""):
    try:
        img_bytes = fig.to_image(format="png", width=1200, height=700, scale=2)
    except Exception:
        st.caption("(Install `kaleido` to enable PPTX slide export for this chart.)")
        return
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.0))
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.text = (fig.layout.title.text or chart_id)
    tf.paragraphs[0].runs[0].font.size = Pt(14)
    tf.paragraphs[0].runs[0].font.bold = True
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    st.download_button("⬇ Download this slide (PPTX)", buf, f"{chart_id}.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"dl_{chart_id}_{key_suffix}")


def render(fig, chart_id, caption="", key_suffix="", height=460):
    style(fig, height=height)
    st.plotly_chart(fig, use_container_width=True, key=f"{chart_id}_{key_suffix}")
    if caption:
        st.caption(caption)
    pptx_button(fig, chart_id, key_suffix)


# ────────────────────────────────────────────────────────────────────────────
# 5b. NARRATIVE HELPERS — plain-English, dynamic text per authority/chart
# ────────────────────────────────────────────────────────────────────────────
def describe_z(z: float) -> str:
    if z >= 1.5:
        return "**substantially above** the national average"
    if z >= 0.5:
        return "above the national average"
    if z <= -1.5:
        return "**substantially below** the national average"
    if z <= -0.5:
        return "below the national average"
    return "close to the national average"


def authority_narrative(name, df_model, Z, selected_indicators, cluster_id, cluster_sizes):
    """Builds a plain-English paragraph describing one authority's standardised profile."""
    z_df = pd.DataFrame(Z, columns=selected_indicators, index=df_model["AREANM"])
    if name not in z_df.index:
        return f"{name} isn't in the current model run (missing a selected indicator)."
    z = z_df.loc[name]
    above = z[z >= 0.5].sort_values(ascending=False)
    below = z[z <= -0.5].sort_values()
    above_txt = ", ".join(INDICATOR_INFO[i]["short"] for i in above.index[:3]) or "none of the selected indicators"
    below_txt = ", ".join(INDICATOR_INFO[i]["short"] for i in below.index[:3]) or "none of the selected indicators"
    size = cluster_sizes.get(cluster_id, 0) if cluster_sizes is not None else None
    txt = (
        f"**{name}** falls into **Cluster {cluster_id}**"
        + (f", alongside {size - 1} other authorities. " if size else ". ")
        + f"Relative to the national average across the selected indicators, {name} is "
          f"notably **higher** on: {above_txt}. It is notably **lower** on: {below_txt}."
    )
    return txt


# ────────────────────────────────────────────────────────────────────────────
# 6. HEADER
# ────────────────────────────────────────────────────────────────────────────
st.markdown(
    f"""
    <div style="background:linear-gradient(90deg,{NAVY},#1E5A8A); padding:1rem 1.4rem;
                border-radius:6px; margin-bottom:1rem;">
    <span style="color:white; font-size:1.05rem; font-weight:bold;">
    🧭 Westminster City Council · Children's Services Directorate</span><br>
    <span style="color:#D9E6F0; font-size:0.9rem;">
    Nearest neighbours &amp; clustering tool, built to support financial planning and benchmarking.</span>
    </div>
    """,
    unsafe_allow_html=True,
)
st.title("Which authorities are genuinely similar to us?")
st.caption(
    "A nearest-neighbours & clustering model built for Children's Services financial planning — "
    "extending the ONS statistical neighbours methodology with ten indicators specific to "
    "children's services demand."
)

c1, c2, c3, c4 = st.columns(4)
c1.metric("Authorities in model", n_areas)
c2.metric("Indicators selected", len(selected_indicators))
c3.metric("Clusters (K)", K, help=f"Best silhouette at K={best_k_silhouette}; elbow knee at K={best_k_elbow}")
c4.metric("Silhouette score", f"{final_sil:.2f}")

tabs = st.tabs([
    "1 · Methodology & caveats",
    "2 · Choosing K (elbow & silhouette)",
    "3 · What do elbow & silhouette mean?",
    "4 · Explore the clusters",
    "5 · Your nearest neighbours",
    "6 · Other validation checks",
])

# ────────────────────────────────────────────────────────────────────────────
# TAB 1 — METHODOLOGY & CAVEATS
# ────────────────────────────────────────────────────────────────────────────
with tabs[0]:
    st.header("Methodology, in plain English")
    st.markdown(
        """
        <div class="method-box">
        <b>What this model does:</b> every authority is described by the indicators you've
        selected in the sidebar. Each indicator is put onto a common scale (so "£ per week" and
        "% in poverty" can be compared fairly), then a computer algorithm (K-means) groups
        authorities that look statistically similar across <i>all</i> of those indicators at once.
        Separately, for the authority you spotlight, we calculate its <b>nearest neighbours</b> —
        the individual authorities whose overall profile is closest to it — which is the number
        that matters most for benchmarking and financial planning.
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.subheader("The pipeline, step by step")
    steps = [
        ("1. Select indicators", "Choose which of the 10 indicators describe your business challenge (sidebar)."),
        ("2. Keep complete cases", "Drop any authority missing even one selected indicator, so every comparison uses the same information."),
        ("3. Cap outliers (winsorise)", "Optionally cap the most extreme 1% of values at each end, so one atypical authority can't dominate the whole model."),
        ("4. Standardise (Z-score)", "Rescale every indicator to mean 0, standard deviation 1 — so a £ measure and a % measure carry equal weight."),
        ("5. Choose K", "Test K = 2 to 9 clusters; use the elbow and silhouette diagnostics (Tab 2) to pick a defensible number."),
        ("6. Cluster (K-means)", "Group authorities that sit close together across all standardised indicators."),
        ("7. Nearest neighbours", "For any one authority, rank every other authority by Euclidean distance across the same indicators."),
    ]
    cols = st.columns(4)
    for i, (title, desc) in enumerate(steps):
        with cols[i % 4]:
            st.markdown(f"**{title}**")
            st.caption(desc)

    st.subheader("Indicators used, and why")
    info_rows = []
    for ind in ALL_INDICATORS:
        meta = INDICATOR_INFO[ind]
        nations = NATION_COVERAGE[ind]
        coverage_str = "UK-wide (all 4 nations)" if nations == ALL_NATIONS else ", ".join(sorted(nations))
        info_rows.append({
            "Indicator": meta["short"], "In current model": "✅" if ind in selected_indicators else "—",
            "What it measures": meta["plain"], "Source": meta["source"], "Period": meta["period"],
            "Geographic coverage": coverage_str, "Vs. ONS method": ONS_METHOD_DIFF[ind],
        })
    st.dataframe(pd.DataFrame(info_rows), use_container_width=True, hide_index=True)

    st.subheader("Which indicators are excluding Scotland and Northern Ireland?")
    missing_nations = ALL_NATIONS - set.intersection(*[NATION_COVERAGE[i] for i in selected_indicators]) \
        if selected_indicators else set()
    if not missing_nations:
        st.success(
            "✅ With your current indicator selection, no nation is structurally excluded — "
            "Scotland and Northern Ireland are both geographically eligible to appear (subject "
            "to ordinary complete-case dropping of individual authorities, not whole nations)."
        )
    else:
        for nation in sorted(missing_nations):
            blockers = [i for i in selected_indicators if nation not in NATION_COVERAGE[i]]
            blocker_txt = ", ".join(INDICATOR_INFO[i]["short"] for i in blockers)
            st.markdown(
                f"""
                <div class="caveat-box">
                <b>{nation} is currently excluded</b> because your selection includes:
                <b>{blocker_txt}</b> — {"this indicator isn't" if len(blockers)==1 else "these indicators aren't"}
                published with a {nation} breakdown. Untick {"it" if len(blockers)==1 else "them"} in the
                sidebar (Children's services indicators) to add {nation} back into the model.
                </div>
                """,
                unsafe_allow_html=True,
            )
        st.caption(
            "Reference (full 10-indicator model): removing **CIN rate + SEN/EHCP rate** restores "
            "Scotland; removing **SEN/EHCP rate + UASC rate** restores Northern Ireland; removing "
            "all three restores both. Adding a nation back doesn't guarantee every one of its "
            "councils survives complete-case dropping — a handful may still be missing on other "
            "grounds (see caveats below)."
        )

    st.subheader("Data caveats — read before presenting this")
    st.markdown(
        f"""
        <div class="caveat-box">
        <b>Different reference periods.</b> Indicators span FYE 2023 to April 2025 (provisional) —
        they are the <i>most recent published</i> figures for each series, not a single snapshot
        year. Census-derived indicators (dependency ratio, population density, no qualifications)
        are fixed at Census 2021 and will not move until Census 2031.<br><br>
        <b>Geographic scope: single-tier authorities only.</b> Children's services indicators
        are published at upper-tier level; some contextual indicators are published at district
        (lower-tier) level for England's ~24 shire counties. Combining the two correctly needs a
        population-weighted district→county aggregation not available in this dataset, so —
        consistent with how ONS/MHCLG themselves restrict statistical-neighbour matching by
        authority type — this model is scoped to <b>{len(single_tier_df)} single-tier authorities</b>
        (unitaries, London boroughs, metropolitan boroughs, and single-tier councils in the
        devolved nations). Two-tier English counties/districts are out of scope for this version.<br><br>
        <b>Nation coverage gaps.</b> Scotland has no direct equivalent to England's "Children in
        Need" statistic, so CIN rate excludes Scotland. Northern Ireland does not publish
        SEN/EHCP or UASC figures broken down by district, so NI is excluded from those two
        indicators. Any authority missing even one selected indicator is dropped from this run
        (currently reducing the model to <b>{n_areas} of {len(single_tier_df)}</b> authorities) —
        this is a real trade-off between breadth of indicators and breadth of geographic coverage.<br><br>
        <b>Deviation from the published ONS method.</b> ONS's own national model does not include
        any children's-social-care-specific indicators (CLA, CP plan, CIN, SEN/EHCP, UASC) — this
        is a bespoke extension using ONS's <i>method</i> (winsorise → standardise → K-means /
        Euclidean nearest neighbours) applied to a children's-services-specific indicator set that
        ONS itself has not modelled.<br><br>
        <b>Child poverty measure basis differs from ONS.</b> This model uses child poverty
        <b>after housing costs (AHC)</b>; ONS's own baseline model uses the measure
        <b>before housing costs (BHC)</b>. The two use the same underlying DWP/HMRC source but are
        not numerically interchangeable — AHC tends to show higher poverty rates in high-rent
        areas (including Westminster) than BHC would for the same households.
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.link_button(
        "ONS statistical neighbours methodology (opens ons.gov.uk)",
        "https://www.ons.gov.uk/peoplepopulationandcommunity/wellbeing/methodologies/"
        "clusteringsimilarlocalauthoritiesandstatisticalnearestneighboursintheukmethodology",
    )

    st.divider()
    st.header("Data cleaning, visualised")
    st.markdown(
        "Every methodology write-up should show, not just tell, how many authorities were lost "
        "at each cleaning step and why. The funnel below does that for your current selection."
    )
    funnel_stages = [
        ("All UK local authority areas in the raw file", int(wide_all["AREACD"].nunique())),
        ("Single-tier authorities only (see caveats above)", len(single_tier_df)),
        (f"Complete cases for your {len(selected_indicators)} selected indicators", n_areas),
    ]
    fig = go.Figure(go.Funnel(
        y=[s[0] for s in funnel_stages], x=[s[1] for s in funnel_stages],
        marker=dict(color=[NAVY, TEAL, ORANGE]),
        textinfo="value+percent initial",
    ))
    fig.update_layout(title="How many authorities survive each cleaning step")
    render(fig, "cleaning_funnel", key_suffix="t1", height=360)
    st.caption(
        f"📊 **In plain English:** starting from {funnel_stages[0][1]} UK local authority areas, "
        f"restricting to single-tier authorities leaves {funnel_stages[1][1]}; requiring every "
        f"selected indicator to be present (no gaps) leaves the {n_areas} authorities the model "
        f"actually runs on. Each step is a genuine, defensible trade-off — see the caveats above "
        f"for exactly why each one happens."
    )

    st.subheader("A real map (optional — add your own boundary file)")
    st.markdown(
        """
        This app runs in a sandboxed cloud environment with no access to Ordnance Survey / ONS
        boundary files, so it cannot draw a real choropleth map on its own — the same reason the
        NI HSC Trust map you built used a locally-downloaded shapefile rather than one fetched
        live. **You can add one in under a minute:**

        1. Download *Local Authority Districts (December 2022) Boundaries UK BUC* as GeoJSON from
           the ONS Open Geography Portal:
           `https://open-geography-portalx-ons.hub.arcgis.com/api/download/v1/items/a2128b32c7fb4205ba99e6344fcbb2be/geojson?layers=0`
        2. Save it into this app's `data/` folder as **`la_boundaries.geojson`**.
        3. Re-run the app — a real choropleth will appear below automatically.
        """
    )
    boundaries = load_boundaries()
    if boundaries is None:
        st.info(
            "No `data/la_boundaries.geojson` found yet, so no live map is shown — the funnel "
            "chart above stands in for it. This is exactly the same shapefile family your NI "
            "Trust map used, just for the whole UK rather than NI alone."
        )
    else:
        map_df = df_model[["AREACD", "AREANM", "Cluster"]].copy()
        map_df["Cluster"] = map_df["Cluster"].astype(str)
        fig = px.choropleth(
            map_df, geojson=boundaries, locations="AREACD", featureidkey=boundaries["_key"],
            color="Cluster", color_discrete_sequence=CLUSTER_COLOURS,
            hover_name="AREANM", title="Your clusters, mapped geographically",
        )
        fig.update_geos(fitbounds="locations", visible=False)
        render(fig, "real_choropleth", key_suffix="t1", height=560)

# ────────────────────────────────────────────────────────────────────────────
# TAB 2 — CHOOSING K
# ────────────────────────────────────────────────────────────────────────────
with tabs[1]:
    st.header("Choosing K: how many groups genuinely exist?")
    st.markdown(
        f"For your current {len(selected_indicators)} selected indicators and {n_areas} authorities: "
        f"the **elbow** sits around **K = {best_k_elbow}**, and the **silhouette score is highest at "
        f"K = {best_k_silhouette}** (score {scores.loc[scores['K']==best_k_silhouette,'Silhouette'].values[0]:.2f}). "
        f"You are currently viewing **K = {K}**."
    )

    left, right = st.columns(2)
    with left:
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=scores["K"], y=scores["Inertia"], mode="lines+markers",
                                  line=dict(color=NAVY, width=3), marker=dict(size=9)))
        annotate_vline(fig, best_k_elbow, f"elbow ≈ K={best_k_elbow}")
        fig.update_layout(title=f"Inertia flattens after K={best_k_elbow} — the 'elbow'",
                           xaxis_title="Number of clusters (K)", yaxis_title="Inertia (lower = tighter clusters)")
        render(fig, "elbow_chart",
               "Inertia = total squared distance from each authority to its cluster's centre. "
               "It always falls as K rises — the elbow is where extra clusters stop buying much.",
               "t2")
    with right:
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=scores["K"], y=scores["Silhouette"], mode="lines+markers",
                                  line=dict(color=TEAL, width=3), marker=dict(size=9)))
        annotate_vline(fig, best_k_silhouette, f"best = K={best_k_silhouette}")
        fig.update_layout(title=f"Cluster separation peaks at K={best_k_silhouette}",
                           xaxis_title="Number of clusters (K)", yaxis_title="Average silhouette score")
        render(fig, "silhouette_chart",
               "Silhouette score (-1 to +1) measures how well-separated the clusters are. "
               "Higher = authorities sit clearly with their own group and far from other groups.",
               "t2")

    st.info(
        f"💡 **Why K={K} was chosen for {len(selected_indicators)} variables:** "
        + (f"the silhouette score is maximised at K={best_k_silhouette}, meaning this is the split "
           f"that produces the most internally-consistent, best-separated groups of authorities — "
           f"adding a 4th, 5th or 6th group starts re-splitting authorities that are already "
           f"well-matched, without improving how distinct the groups are from each other."
           if k_mode.startswith("Auto") else
           f"you have manually overridden the automatic choice (silhouette-optimal K={best_k_silhouette}, "
           f"elbow knee K={best_k_elbow}) to test K={K} directly.")
    )

    st.divider()
    st.header(f"Your {K} clusters: how tight, and how far apart?")
    st.markdown(
        "Elbow and silhouette (above) tell you *which K to pick*. The two charts below apply "
        "that same idea to the **K you've actually chosen**, cluster by cluster — this is the "
        "evidence that a marker (or a finance director) will actually want to see."
    )

    # --- per-cluster tightness (avg within-cluster distance to centroid) ---
    centroids = km_final.cluster_centers_
    within_dist = np.linalg.norm(Z - centroids[km_final.labels_], axis=1)
    tight_df = pd.DataFrame({"Cluster": km_final.labels_, "Dist": within_dist})
    tightness = tight_df.groupby("Cluster")["Dist"].mean().reindex(range(K))

    # --- per-cluster separation (distance from this cluster's centroid to its nearest other centroid) ---
    from scipy.spatial.distance import cdist  # ships with scikit-learn's dependency, safe to use
    centroid_dist = cdist(centroids, centroids)
    np.fill_diagonal(centroid_dist, np.inf)
    separation = pd.Series(centroid_dist.min(axis=1), index=range(K))

    tsep = pd.DataFrame({"Cluster": [f"Cluster {i}" for i in range(K)],
                          "Tightness (lower = tighter)": tightness.values,
                          "Separation (higher = more distinct)": separation.values})

    left2, right2 = st.columns(2)
    with left2:
        fig = px.bar(tsep, x="Cluster", y="Tightness (lower = tighter)", color="Cluster",
                     color_discrete_sequence=CLUSTER_COLOURS,
                     title="How tight is each cluster internally?")
        fig.update_layout(showlegend=False, xaxis_title="", yaxis_title="Avg. distance to cluster centre")
        fig.update_yaxes(rangemode="tozero")
        render(fig, "tightness_chart",
               "A short bar means that cluster's authorities sit close together — a genuinely "
               "coherent comparator group. A tall bar means the cluster is more loosely defined.",
               "t2")
    with right2:
        fig = px.bar(tsep, x="Cluster", y="Separation (higher = more distinct)", color="Cluster",
                     color_discrete_sequence=CLUSTER_COLOURS,
                     title="How distinct is each cluster from its nearest neighbouring cluster?")
        fig.update_layout(showlegend=False, xaxis_title="", yaxis_title="Distance to nearest other cluster centre")
        fig.update_yaxes(rangemode="tozero")
        render(fig, "separation_chart",
               "A tall bar means that cluster is clearly its own group. A short bar means it "
               "sits close to — and could plausibly be merged with — a neighbouring cluster.",
               "t2")

    tightest = tsep.loc[tsep["Tightness (lower = tighter)"].idxmin(), "Cluster"]
    loosest = tsep.loc[tsep["Tightness (lower = tighter)"].idxmax(), "Cluster"]
    most_distinct = tsep.loc[tsep["Separation (higher = more distinct)"].idxmax(), "Cluster"]
    st.info(
        f"💡 **Reading these two charts together:** {tightest} is the most internally consistent "
        f"group (smallest spread), while {loosest} is the loosest — its members are more varied. "
        f"{most_distinct} stands furthest apart from its nearest neighbouring cluster, making it "
        f"the most clearly distinct group in this model."
    )

    st.subheader(f"Every authority's individual fit, K={K} (real data)")
    order = np.argsort(df_model["Cluster"].values * 1000 - sample_sil)
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=sample_sil[order], y=df_model["AREANM"].values[order], orientation="h",
        marker_color=[CLUSTER_COLOURS[c % len(CLUSTER_COLOURS)] for c in df_model["Cluster"].values[order]],
    ))
    fig.update_layout(title=f"Silhouette width for every one of your {n_areas} authorities",
                       xaxis_title="Silhouette width (-1 to +1)", yaxis_title="")
    fig.update_yaxes(showticklabels=False)
    annotate_vline(fig, final_sil, f"model average = {final_sil:.2f}", color=DARK_GREY)
    render(fig, "real_silhouette_diagram",
           "Each bar is one authority, grouped and coloured by cluster, sorted from best-fitting "
           "to worst within each group. Bars near zero or negative sit on a cluster boundary — "
           "worth a sanity check before quoting them as a clean comparator.",
           "t2", height=520)

# ────────────────────────────────────────────────────────────────────────────
# TAB 3 — DEDICATED ELBOW / SILHOUETTE EXPLAINER (toy example, plain English)
# ────────────────────────────────────────────────────────────────────────────
with tabs[2]:
    st.header("What do 'elbow' and 'silhouette' actually mean?")
    st.caption("A worked example with 12 illustrative points, before we apply either idea to real authorities.")

    rng = np.random.default_rng(7)
    toy = np.vstack([
        rng.normal([2, 2], 0.5, (4, 2)),
        rng.normal([8, 2], 0.5, (4, 2)),
        rng.normal([5, 7], 0.5, (4, 2)),
    ])
    toy_df = pd.DataFrame(toy, columns=["x", "y"])

    ex1, ex2 = st.columns(2)
    with ex1:
        st.subheader("The elbow method (inertia)")
        st.markdown(
            """
            **In plain English:** inertia measures how tightly points sit around their cluster's
            centre — like measuring how compact each group is. As you add more clusters, inertia
            *always* falls (more groups = shorter distances to the nearest centre). The **elbow**
            is where the curve stops dropping steeply and starts flattening out — meaning extra
            clusters aren't finding genuinely new structure, they're just splitting groups that
            were already fine.
            """
        )
        toy_inertias = []
        for k in range(1, 7):
            km = KMeans(n_clusters=k, random_state=1, n_init=10).fit(toy)
            toy_inertias.append(km.inertia_)
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=list(range(1, 7)), y=toy_inertias, mode="lines+markers",
                                  line=dict(color=NAVY, width=3), marker=dict(size=10)))
        annotate_vline(fig, 3, "elbow at K=3")
        fig.update_layout(title="Toy example: the elbow sits at K=3 (the 3 true groups)",
                           xaxis_title="K", yaxis_title="Inertia")
        render(fig, "toy_elbow", key_suffix="t3")

    with ex2:
        st.subheader("The silhouette score")
        st.markdown(
            """
            **In plain English:** for every single point, silhouette compares (a) how close it is
            to *other points in its own cluster* against (b) how close it is to the *nearest
            other cluster*. A score near **+1** means "clearly belongs here, far from anywhere
            else." A score near **0** means "sitting right on the boundary between two groups."
            A **negative** score means a point may actually be in the wrong cluster. Averaging
            this across every point gives one overall score for a given K.
            """
        )
        km3 = KMeans(n_clusters=3, random_state=1, n_init=10).fit(toy)
        sil_vals = silhouette_samples(toy, km3.labels_)
        order = np.argsort(km3.labels_ * 100 - sil_vals)
        fig = go.Figure()
        colours_map = {0: NAVY, 1: ORANGE, 2: TEAL}
        fig.add_trace(go.Bar(
            x=sil_vals[order], y=[f"pt {i}" for i in order], orientation="h",
            marker_color=[colours_map[c] for c in km3.labels_[order]],
        ))
        annotate_vline(fig, sil_vals.mean(), f"average = {sil_vals.mean():.2f}", color=DARK_GREY)
        fig.update_layout(title="Every point's individual silhouette width (K=3)",
                           xaxis_title="Silhouette width (-1 to +1)", yaxis_title="")
        render(fig, "toy_silhouette", key_suffix="t3")

    st.subheader("See the 3 toy groups the scores above are describing")
    km3 = KMeans(n_clusters=3, random_state=1, n_init=10).fit(toy)
    toy_df["Cluster"] = km3.labels_.astype(str)
    fig = px.scatter(toy_df, x="x", y="y", color="Cluster",
                      color_discrete_sequence=[NAVY, ORANGE, TEAL],
                      title="The three groups these 12 points naturally form")
    fig.update_traces(marker=dict(size=16, line=dict(width=1, color="white")))
    render(fig, "toy_scatter", key_suffix="t3")

    st.markdown(
        """
        <div class="method-box">
        <b>Why we use both together:</b> the elbow method is quick but the "bend" can be a
        matter of judgement. The silhouette score is more precise but only ever gives you one
        K at a time. Using both — and picking the K where they agree, or explaining clearly when
        they don't — is standard practice and is exactly what ONS do in their own published
        methodology.
        </div>
        """,
        unsafe_allow_html=True,
    )

# ────────────────────────────────────────────────────────────────────────────
# TAB 4 — EXPLORE THE CLUSTERS (real data)
# ────────────────────────────────────────────────────────────────────────────
with tabs[3]:
    st.header(f"The {K} groups found in your {n_areas} authorities")

    cluster_sizes = df_model["Cluster"].value_counts().sort_index()
    fig = px.bar(x=[f"Cluster {i}" for i in cluster_sizes.index], y=cluster_sizes.values,
                 color=[f"Cluster {i}" for i in cluster_sizes.index],
                 color_discrete_sequence=CLUSTER_COLOURS,
                 title="How many authorities fall into each group")
    fig.update_layout(showlegend=False, xaxis_title="", yaxis_title="Number of authorities")
    fig.update_yaxes(rangemode="tozero")
    render(fig, "cluster_sizes", key_suffix="t4")
    st.caption(
        f"📊 **In plain English:** your {n_areas} authorities split into {K} groups ranging from "
        f"{cluster_sizes.min()} to {cluster_sizes.max()} authorities each. Uneven cluster sizes "
        f"aren't a problem in themselves — some genuine 'types' of authority (e.g. inner-London "
        f"boroughs) are simply rarer nationally than others."
    )

    spot_row = df_model[df_model["AREANM"] == spotlight]
    spot_cluster = int(spot_row["Cluster"].iloc[0]) if len(spot_row) else None

    with st.expander("ℹ️ What is PCA, and why does the next chart use it?", expanded=False):
        st.markdown(
            f"""
            Your model uses **{len(selected_indicators)} indicators** at once — impossible to plot
            directly (that would need {len(selected_indicators)} axes). **PCA (Principal Component
            Analysis)** solves this by mathematically combining all {len(selected_indicators)}
            indicators into just 2 new "super-axes" that capture as much of the original variation
            as possible, so every authority can be placed on one simple 2D map.

            **Important:** PCA is used *only* to draw the chart below — the clustering itself was
            already done on all {len(selected_indicators)} original indicators (Tab 2). Two
            authorities that look close on this map are close across their *whole* profile, not
            just on two indicators.

            For your current selection, these 2 axes capture **{explained.sum()*100:.0f}%** of the
            total variation across all {len(selected_indicators)} indicators — the closer to 100%,
            the more trustworthy this 2D picture is as a stand-in for the full model.
            """
        )
        fig = px.bar(x=["Axis 1 (PC1)", "Axis 2 (PC2)"], y=explained * 100,
                     title="How much of the original variation each PCA axis captures",
                     color_discrete_sequence=[NAVY])
        fig.update_layout(showlegend=False, xaxis_title="", yaxis_title="% of variation explained")
        fig.update_yaxes(rangemode="tozero")
        render(fig, "pca_variance", key_suffix="t4", height=340)

    left, right = st.columns([3, 2])
    with left:
        fig = px.scatter(
            df_model[df_model["AREANM"] != spotlight], x="PC1", y="PC2",
            color=df_model[df_model["AREANM"] != spotlight]["Cluster"].astype(str),
            hover_name="AREANM", color_discrete_sequence=CLUSTER_COLOURS,
            title=f"Every authority's overall profile, compressed to 2 dimensions "
                  f"({explained.sum()*100:.0f}% of variation captured)",
            labels={"color": "Cluster"},
        )
        fig.update_traces(marker=dict(size=10, line=dict(width=0.5, color="white")))
        if len(spot_row):
            fig.add_trace(go.Scatter(
                x=spot_row["PC1"], y=spot_row["PC2"], mode="markers+text",
                marker=dict(size=20, symbol="star", color=ORANGE, line=dict(width=1.5, color=DARK_GREY)),
                text=[spotlight], textposition="top center", textfont=dict(size=12, color=DARK_GREY),
                name=spotlight, showlegend=True,
            ))
        render(fig, "pca_scatter",
               "This chart (PCA) squashes every selected indicator down to 2 axes so similar "
               "authorities appear near each other — it's a map of similarity, not geography.",
               "t4")
    with right:
        st.subheader(f"{spotlight}'s profile")
        if len(spot_row):
            st.metric("Assigned to", f"Cluster {spot_cluster}")
            st.metric("Silhouette width", f"{spot_row['SilhouetteWidth'].iloc[0]:.2f}",
                       help="How clearly this authority belongs to its group vs the next-nearest group.")
            st.caption(f"{cluster_sizes[spot_cluster]} authorities share this cluster.")
        else:
            st.warning(f"{spotlight} is missing one or more selected indicators, so it was "
                        f"dropped from this run (see Tab 1 caveats).")

    st.markdown(authority_narrative(spotlight, df_model, Z, selected_indicators, spot_cluster, cluster_sizes))

    st.divider()
    st.subheader("Cluster profiles — what makes each group distinct")
    profile = df_model.groupby("Cluster")[selected_indicators].mean()
    profile_z = (profile - df_model[selected_indicators].mean()) / df_model[selected_indicators].std()
    heat = profile_z.rename(columns={k: INDICATOR_INFO[k]["short"] for k in selected_indicators})
    fig = px.imshow(
        heat.T, color_continuous_scale=[[0, TEAL], [0.5, "white"], [1, ORANGE]],
        aspect="auto", labels=dict(color="Std. dev. from national average"),
        title="Where each cluster sits above (orange) or below (teal) the national average",
    )
    fig.update_xaxes(title="Cluster", tickvals=list(range(K)), ticktext=[f"Cluster {i}" for i in range(K)])
    fig.update_yaxes(title="")
    render(fig, "cluster_heatmap",
           "Each cell shows how far that cluster's average sits from the overall average, in "
           "standard deviations — the same scale the clustering itself used.", "t4", height=380 + 18 * len(selected_indicators))

    if spot_cluster is not None:
        cluster_profile_row = profile_z.loc[spot_cluster]
        top_defining = cluster_profile_row.abs().sort_values(ascending=False).index[:3]
        defining_txt = "; ".join(
            f"**{INDICATOR_INFO[i]['short']}** ({describe_z(cluster_profile_row[i])})" for i in top_defining
        )
        st.caption(f"📊 **Cluster {spot_cluster}'s defining characteristics:** {defining_txt}.")

    st.subheader("All selected indicators at once, by cluster")
    pc_df = df_model[["AREANM", "Cluster"] + selected_indicators].copy()
    dims = [dict(label=INDICATOR_INFO[c]["short"], values=pc_df[c]) for c in selected_indicators]
    fig = go.Figure(data=go.Parcoords(
        line=dict(color=pc_df["Cluster"], colorscale=[[i/(max(K-1,1)), c] for i, c in enumerate(CLUSTER_COLOURS[:K])]),
        dimensions=dims,
    ))
    fig.update_layout(title="Each line is one authority — coloured by cluster",
                       margin=dict(l=80, r=80, t=90, b=40))
    render(fig, "parallel_coords",
           "Parallel coordinates: authorities in the same cluster tend to trace similar paths "
           "across all indicators at once, showing which measures separate the groups most clearly.",
           "t4", height=480)

    if spotlight in df_model["AREANM"].values:
        st.divider()
        st.subheader(f"{spotlight} vs its own cluster vs the national average")
        z_df = pd.DataFrame(Z, columns=selected_indicators, index=df_model["AREANM"])
        spot_z = z_df.loc[spotlight]
        cluster_avg_z = z_df[df_model["Cluster"].values == spot_cluster].mean()
        radar_labels = [INDICATOR_INFO[c]["short"] for c in selected_indicators]
        fig = go.Figure()
        fig.add_trace(go.Scatterpolar(r=list(spot_z.values) + [spot_z.values[0]],
                                       theta=radar_labels + [radar_labels[0]],
                                       name=spotlight, line=dict(color=ORANGE, width=3)))
        fig.add_trace(go.Scatterpolar(r=list(cluster_avg_z.values) + [cluster_avg_z.values[0]],
                                       theta=radar_labels + [radar_labels[0]],
                                       name=f"Cluster {spot_cluster} average", line=dict(color=NAVY, width=2, dash="dash")))
        fig.add_trace(go.Scatterpolar(r=[0] * (len(radar_labels) + 1), theta=radar_labels + [radar_labels[0]],
                                       name="National average", line=dict(color=GREY, width=1, dash="dot")))
        fig.update_layout(title=f"{spotlight} (standard deviations from the national average)",
                           polar=dict(radialaxis=dict(visible=True),
                                       angularaxis=dict(tickfont=dict(size=10))))
        render(fig, "radar_chart", key_suffix="t4", height=560)
        biggest_gap = (spot_z - cluster_avg_z).abs().sort_values(ascending=False).index[0]
        st.caption(
            f"📊 **Reading this chart:** where the orange line sits outside the navy dashed line, "
            f"{spotlight} is more extreme than even its own cluster's average. The largest gap "
            f"between {spotlight} and its cluster is on "
            f"**{INDICATOR_INFO[biggest_gap]['short']}** — worth flagging explicitly if you "
            f"present this cluster as a comparator group."
        )

    st.divider()
    st.subheader("How do two indicators relate to each other? (exploratory)")
    st.caption(
        "A simple linear regression between any two of your selected indicators — useful "
        "supporting evidence for *why* certain indicators move together, separate from the "
        "clustering model itself."
    )
    rc1, rc2 = st.columns(2)
    x_var = rc1.selectbox("X axis", selected_indicators, index=0,
                            format_func=lambda x: INDICATOR_INFO[x]["short"], key="reg_x")
    y_default = 1 if len(selected_indicators) > 1 else 0
    y_var = rc2.selectbox("Y axis", selected_indicators, index=y_default,
                            format_func=lambda x: INDICATOR_INFO[x]["short"], key="reg_y")
    if x_var == y_var:
        st.warning("Pick two different indicators to see a relationship.")
    else:
        xv = df_model[x_var].values.astype(float)
        yv = df_model[y_var].values.astype(float)
        slope, intercept = np.polyfit(xv, yv, 1)
        pred = slope * xv + intercept
        ss_res = np.sum((yv - pred) ** 2)
        ss_tot = np.sum((yv - yv.mean()) ** 2)
        r2 = 1 - ss_res / ss_tot if ss_tot > 0 else 0
        corr = np.corrcoef(xv, yv)[0, 1]
        line_x = np.linspace(xv.min(), xv.max(), 50)
        line_y = slope * line_x + intercept

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=xv, y=yv, mode="markers", name="Authorities",
                                  marker=dict(size=8, color=GREY, line=dict(width=0.5, color="white")),
                                  text=df_model["AREANM"], hovertemplate="%{text}<extra></extra>"))
        if len(spot_row):
            fig.add_trace(go.Scatter(x=spot_row[x_var], y=spot_row[y_var], mode="markers+text",
                                      marker=dict(size=16, symbol="star", color=ORANGE),
                                      text=[spotlight], textposition="top center", name=spotlight))
        fig.add_trace(go.Scatter(x=line_x, y=line_y, mode="lines",
                                  line=dict(color=NAVY, width=3), name="Fitted trend line"))
        fig.update_layout(
            title=f"{INDICATOR_INFO[y_var]['short']} vs {INDICATOR_INFO[x_var]['short']} "
                  f"(R² = {r2:.2f}, correlation = {corr:.2f})",
            xaxis_title=INDICATOR_INFO[x_var]["short"], yaxis_title=INDICATOR_INFO[y_var]["short"],
        )
        render(fig, "regression_chart", key_suffix="t4")

        strength = ("a strong" if abs(corr) >= 0.7 else "a moderate" if abs(corr) >= 0.4 else "a weak")
        direction = "positive" if corr > 0 else "negative"
        st.caption(
            f"📊 **In plain English:** there is {strength} {direction} relationship between "
            f"{INDICATOR_INFO[x_var]['short'].lower()} and {INDICATOR_INFO[y_var]['short'].lower()} "
            f"across these {n_areas} authorities (R² = {r2:.2f} means the trend line explains "
            f"{r2*100:.0f}% of the variation in {INDICATOR_INFO[y_var]['short'].lower()}). "
            f"This is exploratory evidence for *why* two indicators may move together — it is not "
            f"the clustering model itself, and correlation here does not prove one causes the other."
        )

# ────────────────────────────────────────────────────────────────────────────
# TAB 5 — NEAREST NEIGHBOURS
# ────────────────────────────────────────────────────────────────────────────
with tabs[4]:
    st.header(f"{spotlight}'s nearest neighbours")
    if spotlight not in df_model["AREANM"].values:
        st.warning(f"{spotlight} was dropped from this run because it is missing one of the "
                   f"selected indicators — remove that indicator in the sidebar, or pick a "
                   f"different spotlight authority.")
    else:
        st.markdown(
            """
            <div class="method-box">
            <b>This is separate from clustering.</b> Clustering sorts every authority into a
            small number of broad groups. Nearest neighbours instead asks one precise question:
            <i>of every other authority, which ones is this one most similar to, overall?</i>
            It's calculated as the straight-line (Euclidean) distance between standardised
            indicator profiles — the smaller the distance, the more similar the two authorities.
            </div>
            """,
            unsafe_allow_html=True,
        )

        z_df = pd.DataFrame(Z, columns=selected_indicators, index=df_model["AREANM"])
        target = z_df.loc[spotlight]
        dist = np.sqrt(((z_df - target) ** 2).sum(axis=1))
        dist = dist.drop(spotlight).sort_values()
        n_neighbours = st.slider("How many neighbours to show", 5, 20, 10)
        top = dist.head(n_neighbours).reset_index()
        top.columns = ["Authority", "Distance"]
        top["Tier"] = top["Authority"].map(df_model.set_index("AREANM")["Tier"])

        left, right = st.columns([3, 2])
        with left:
            fig = px.bar(top.sort_values("Distance", ascending=True), x="Distance", y="Authority",
                         orientation="h", color="Distance", color_continuous_scale=[[0, ORANGE], [1, GREY]],
                         title=f"Closest matches to {spotlight} (shortest bar = most similar)")
            fig.update_yaxes(categoryorder="total descending")
            fig.update_layout(coloraxis_showscale=False)
            render(fig, "neighbours_bar", key_suffix="t5")
        with right:
            median_national = dist.median()
            st.metric(f"{spotlight}'s nearest neighbour distance", f"{dist.min():.2f}")
            st.metric("Median distance across all pairs in this model", f"{median_national:.2f}")
            if dist.min() > median_national:
                st.caption(
                    f"⚠️ {spotlight}'s closest match is *further away* than the typical "
                    f"authority's closest match — a sign {spotlight} is a statistical outlier "
                    f"and comparisons should be made cautiously."
                )
            else:
                st.caption(f"{spotlight} has at least one closely-matched comparator authority.")

        st.subheader("Raw indicator values — for the finance conversation")
        table_cols = ["AREANM"] + selected_indicators
        display = df_model[df_model["AREANM"].isin([spotlight] + top["Authority"].tolist())][table_cols]
        display = display.set_index("AREANM").loc[[spotlight] + top["Authority"].tolist()]
        display = display.rename(columns={k: INDICATOR_INFO[k]["short"] for k in selected_indicators})
        st.dataframe(display.style.format("{:.1f}").apply(
            lambda s: ["background-color:#FDEEE3" if s.name == spotlight else "" for _ in s], axis=1),
            use_container_width=True)
        st.caption(
            "Unlike the charts above (which use standardised/scaled values so different units "
            "compare fairly), this table shows the real published figures — the numbers to "
            "actually quote in a finance report."
        )

        csv = display.reset_index().to_csv(index=False).encode()
        st.download_button("⬇ Download this table (CSV)", csv, f"{spotlight}_nearest_neighbours.csv", "text/csv")

# ────────────────────────────────────────────────────────────────────────────
# TAB 6 — OTHER CLUSTER VALIDATION CHECKS
# ────────────────────────────────────────────────────────────────────────────
with tabs[5]:
    st.header("Beyond elbow & silhouette: three more checks")
    st.markdown(
        """
        <div class="method-box">
        Elbow and silhouette are the two best-known cluster diagnostics, but relying on either
        alone has known limitations: the elbow's "bend" is often genuinely ambiguous to read by
        eye, and silhouette can be pulled around by clusters of very different sizes or densities.
        Good practice — and the theme of every methods-comparison writeup on this topic — is to
        triangulate: check whether the data has real cluster structure at all, cross-check K
        against indices built on a different mathematical basis, and test whether the clusters
        hold up under resampling. None of this replaces elbow/silhouette; it stress-tests them.
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.subheader("① Is this data even clusterable? (Hopkins statistic)")
    st.markdown(
        "Before asking *how many* clusters, it's worth asking whether there is genuine group "
        "structure here at all, rather than one smooth, structureless cloud of authorities that "
        "K-means would carve up arbitrarily regardless. The **Hopkins statistic** compares how "
        "close real authorities sit to their nearest neighbour against how close a set of "
        "randomly-scattered points would sit to *their* nearest real authority."
    )
    H = hopkins_statistic(Z)
    hop_col1, hop_col2 = st.columns([2, 3])
    with hop_col1:
        fig = go.Figure(go.Indicator(
            mode="gauge+number", value=H, number=dict(valueformat=".2f"),
            gauge=dict(
                axis=dict(range=[0, 1]),
                bar=dict(color=ORANGE if H > 0.75 else NAVY),
                steps=[
                    dict(range=[0, 0.5], color="#F0F0F0"),
                    dict(range=[0.5, 0.75], color="#D9E6F0"),
                    dict(range=[0.75, 1], color="#CFE8E4"),
                ],
                threshold=dict(line=dict(color=DARK_GREY, width=3), value=0.5),
            ),
            title=dict(text="Hopkins statistic"),
        ))
        render(fig, "hopkins_gauge", key_suffix="t6", height=320)
    with hop_col2:
        verdict = (
            "well above the 0.5 'random data' line — this data has genuine, meaningful cluster "
            "structure, so K-means clustering is a legitimate approach here (not just imposing "
            "arbitrary groups on noise)."
            if H > 0.7 else
            "close to the 0.5 line that indicates essentially random, structureless data — "
            "treat any clusters found on this selection with real caution, and consider "
            "whether the chosen indicators genuinely separate authorities into distinct types."
        )
        st.markdown(
            f"**Your score: {H:.2f}.** A score of 0.5 means no more cluster structure than "
            f"randomly scattered points would show; a score approaching 1.0 means strong, "
            f"genuine grouping. Your current indicator selection sits {verdict}"
        )

    st.divider()
    st.subheader("② Do independent indices agree on K?")
    st.markdown(
        "Silhouette is one way to score a given K. **Calinski-Harabasz** (the ratio of "
        "between-cluster to within-cluster spread — higher is better) and **Davies-Bouldin** "
        "(average similarity between each cluster and its most-similar neighbour — lower is "
        "better) are built on different maths entirely. When several independent indices point "
        "to a similar K, that's much stronger evidence than any one of them alone."
    )
    other_scores = other_validation_scores(Z)
    best_k_ch = int(other_scores.loc[other_scores["Calinski-Harabasz"].idxmax(), "K"])
    best_k_db = int(other_scores.loc[other_scores["Davies-Bouldin"].idxmin(), "K"])

    vc1, vc2 = st.columns(2)
    with vc1:
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=other_scores["K"], y=other_scores["Calinski-Harabasz"],
                                  mode="lines+markers", line=dict(color=NAVY, width=3), marker=dict(size=9)))
        annotate_vline(fig, best_k_ch, f"best = K={best_k_ch}")
        fig.update_layout(title="Calinski-Harabasz index (higher = better-separated clusters)",
                           xaxis_title="K", yaxis_title="Calinski-Harabasz score")
        render(fig, "ch_chart", key_suffix="t6")
    with vc2:
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=other_scores["K"], y=other_scores["Davies-Bouldin"],
                                  mode="lines+markers", line=dict(color=TEAL, width=3), marker=dict(size=9)))
        annotate_vline(fig, best_k_db, f"best = K={best_k_db}")
        fig.update_layout(title="Davies-Bouldin index (lower = better-separated clusters)",
                           xaxis_title="K", yaxis_title="Davies-Bouldin score")
        render(fig, "db_chart", key_suffix="t6")

    agree = len({best_k_silhouette, best_k_ch, best_k_db}) == 1
    votes_df = pd.DataFrame({
        "Method": ["Elbow (inertia)", "Silhouette", "Calinski-Harabasz", "Davies-Bouldin"],
        "Suggests K": [best_k_elbow, best_k_silhouette, best_k_ch, best_k_db],
    })
    st.dataframe(votes_df, use_container_width=True, hide_index=True)
    if agree:
        st.success(
            f"✅ Silhouette, Calinski-Harabasz and Davies-Bouldin all independently point to "
            f"**K={best_k_silhouette}** — strong, cross-validated evidence for this choice."
        )
    else:
        st.warning(
            f"⚠️ The indices don't fully agree (K={best_k_elbow}/{best_k_silhouette}/{best_k_ch}/"
            f"{best_k_db} across elbow/silhouette/Calinski-Harabasz/Davies-Bouldin). This isn't "
            f"unusual and doesn't invalidate the model — but it's worth naming explicitly rather "
            f"than quoting a single K as if it were unambiguous. Report the range and explain why "
            f"you ultimately chose K={K}."
        )

    st.divider()
    st.subheader("③ Do the clusters survive resampling? (bootstrap stability)")
    st.markdown(
        f"If K={K} reflects real structure, clustering a random 80% subsample of authorities "
        f"should recover *almost* the same groups as the full model. If it reflects noise or one "
        f"or two unusual authorities, the groups will shuffle unpredictably each time. This runs "
        f"K-means on 30 different random subsamples and compares each to the full model using the "
        f"**Adjusted Rand Index** (1.0 = identical groupings, 0.0 = no better than chance)."
    )
    aris = bootstrap_stability(Z, K)
    fig = go.Figure()
    fig.add_trace(go.Histogram(x=aris, nbinsx=15, marker_color=NAVY, opacity=0.85))
    annotate_vline(fig, float(np.mean(aris)), f"mean = {np.mean(aris):.2f}")
    fig.update_layout(title=f"Cluster stability across 30 resamples (K={K})",
                       xaxis_title="Adjusted Rand Index vs. full model", yaxis_title="Number of resamples")
    render(fig, "stability_hist", key_suffix="t6")

    mean_ari = float(np.mean(aris))
    if mean_ari >= 0.75:
        stab_msg = "highly stable — the clusters are a robust feature of the data, not a fluke of exactly which authorities happen to be in this run."
    elif mean_ari >= 0.5:
        stab_msg = "moderately stable — the broad grouping holds up, but individual authorities near cluster boundaries may switch groups on slightly different data."
    else:
        stab_msg = "not very stable — treat cluster membership as indicative rather than definitive, and lean more heavily on the individual nearest-neighbours results (Tab 5) than on cluster labels for financial decisions."
    st.info(f"💡 **Mean ARI = {mean_ari:.2f}.** This clustering is {stab_msg}")

st.divider()
st.caption(
    "Built as a DAP (Data Analytics Principles) artefact. Methodology extends the ONS "
    "statistical nearest neighbours approach with children's-services-specific indicators. "
    "See Tab 1 for full data sources, periods and caveats."
)
